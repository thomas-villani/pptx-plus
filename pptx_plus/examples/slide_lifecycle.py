"""The three v0.1 verbs, end to end, on a deck built from scratch.

    python -m pptx_plus.examples.slide_lifecycle [--out DIR]

Builds a five-slide deck with a picture and a chart, then reorders, copies and
removes slides -- printing the running order after each step and checking the
saved package against the integrity battery, so the example doubles as a
demonstration that the result is actually well formed.

Writes to a temp directory unless --out says otherwise. It will not write into
your working directory uninvited.
"""

from __future__ import annotations

import argparse
import io
import tempfile
from pathlib import Path

import pptx
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

# `pptx.Presentation` is a FACTORY FUNCTION, not the class. The class is
# `pptx.presentation.Presentation`, and it is what a type annotation needs --
# annotating with the factory is an error under `mypy --strict`. Hence
# `pptx.Presentation()` to build one and this import to name the type.
from pptx.presentation import Presentation
from pptx.util import Inches

from pptx_plus._testing import assert_package_integrity, saved
from pptx_plus.slides import contains, delete_slide, duplicate_slide, move_slide

# A 1x1 GIF. The example is about the slide lifecycle, not about the picture.
GIF = (
    b"GIF89a\x01\x00\x01\x00\x80\x00\x00\x00\x00\x00\xff\xff\xff!"
    b"\xf9\x04\x01\x00\x00\x00\x00,\x00\x00\x00\x00\x01\x00\x01\x00\x00\x02\x02D\x01\x00;"
)
TITLE_ONLY = 5


def build_deck() -> Presentation:
    prs = pptx.Presentation()
    for name in ("Agenda", "Market", "Results", "Roadmap", "Questions"):
        slide = prs.slides.add_slide(prs.slide_layouts[TITLE_ONLY])
        slide.shapes.title.text = name

    prs.slides[1].shapes.add_picture(io.BytesIO(GIF), Inches(1), Inches(2))

    data = CategoryChartData()
    data.categories = ["Q1", "Q2", "Q3"]
    data.add_series("Revenue", (10.0, 20.0, 30.0))
    prs.slides[2].shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(2), Inches(6), Inches(4), data
    )
    return prs


def order(prs: Presentation) -> str:
    return ", ".join(slide.shapes.title.text for slide in prs.slides)


def show(prs: Presentation, label: str) -> None:
    """Print the running order and confirm the saved package is still valid."""
    assert_package_integrity(saved(prs))
    print(f"  {label:<34} {order(prs)}")


def main() -> None:
    parser = argparse.ArgumentParser(description="pptx_plus slide lifecycle demo")
    parser.add_argument("--out", type=Path, default=None, help="output directory")
    args = parser.parse_args()
    out_dir = args.out or Path(tempfile.mkdtemp(prefix="pptx-plus-example-"))
    out_dir.mkdir(parents=True, exist_ok=True)

    prs = build_deck()
    print("Slide lifecycle -- each step verified against the saved package.\n")
    show(prs, "built")

    move_slide(prs, 4, 1)
    show(prs, "move_slide(4, 1)")

    copy = duplicate_slide(prs, 3)
    show(prs, "duplicate_slide(3)")
    print(f"\n  The copy carries the chart, and owns it: slide {prs.slides.index(copy)}.")
    print("  Its embedded workbook is a separate part, so editing the copy")
    print("  cannot reach back into the original.\n")

    delete_slide(prs, 0)
    show(prs, "delete_slide(0)")

    doomed = prs.slides[0]
    delete_slide(prs, doomed)
    show(prs, "delete_slide(slide)")

    # A deleted Slide stays alive and readable -- deletion detaches its part
    # from the relationship graph, it destroys nothing. SPEC 5.7.
    print(f'\n  The deleted slide object is still readable: "{doomed.shapes.title.text}"')
    print(f"  but it is no longer part of the deck: contains() -> {contains(prs, doomed)}")

    path = out_dir / "slide-lifecycle.pptx"
    # `str()` because python-pptx annotates `save` as `str | IO[bytes]`; a
    # Path works at runtime but not under `mypy --strict`.
    prs.save(str(path))
    print(f"\nSaved {path} ({path.stat().st_size:,} bytes)")


if __name__ == "__main__":
    main()

"""Duplicate a slide within its deck.

python-pptx's other twelve-year-old feature request, and the one where the
circulating recipe is not merely incomplete but actively wrong.

``copy.deepcopy(slide.element)`` copies the ``<p:spTree>`` faithfully and
carries every ``r:id`` and ``r:embed`` inside it across unchanged. Those ids
are scoped to the part that contains them (SPEC §3.2), so on the new part they
resolve against a different relationship set — where they name something else,
or nothing. The failure is silent: no exception, no repair prompt, just a
picture that does not appear.

Doing it correctly means cloning the part graph, re-minting every
relationship, and rewriting every relationship id in the copied XML. All three
live in ``core``; this module is the verb that composes them with the deck's
running order, its slide-id allocation, and its sections.

ECMA-376 Part 1 §19.2.1.34 (``p:sldIdLst``), §13.2 (part reachability).

SPEC §5.4. This module imports only from ``pptx_plus.core`` (SPEC §9.1).
"""

from __future__ import annotations

from typing import TYPE_CHECKING, cast

from pptx.opc.constants import RELATIONSHIP_TYPE as RT

from pptx_plus.core.clone import ClonePolicy, clone_part_graph
from pptx_plus.core.oxml import sld_id_lst
from pptx_plus.core.sections import insert_slide
from pptx_plus.slides.resolve import SlideIndexError, resolve_slide

if TYPE_CHECKING:
    from pptx.presentation import Presentation
    from pptx.slide import Slide


def duplicate_slide(
    prs: Presentation,
    slide_or_index: Slide | int,
    *,
    to_index: int | None = None,
    with_notes: bool = True,
) -> Slide:
    """Duplicate a slide and return the copy.

    Args:
        prs: The presentation to duplicate within.
        slide_or_index: The slide to copy, or its index. Negative indices
            count from the end.
        to_index: Where the copy goes, as a position in the **resulting**
            deck. Defaults to immediately after the source. Negative values
            count from the end; ``to_index=-1`` appends.
        with_notes: Copy the slide's speaker notes.

    Returns:
        The new :class:`~pptx.slide.Slide`.

    Raises:
        SlideIndexError: Either index is outside the deck.
        SlideNotFoundError: The slide given is not in this deck.
        TypeError: ``slide_or_index`` is neither a ``Slide`` nor an ``int``.
        DanglingRelationshipError: The source slide's XML names a relationship
            that does not exist. Real decks do contain these, and copying one
            into a slide this library produced would be worse than refusing.

    **What is shared and what is copied.** Images and media are shared by
    reference: an image part is identified by its bytes, and two slides
    pointing at one is the format's own way of saying "the same picture."
    Charts with their embedded workbooks, SmartArt definition parts, embedded
    objects and — unless ``with_notes=False`` — the notes slide are copied, so
    editing the duplicate cannot reach back into the original. Layouts,
    masters and themes are reused, because two slides sharing a layout is
    correct and duplicating the layout would not be.

    A slide-jump hyperlink points at the *same* slide it did before. The copy
    links to slide 4; it does not bring a second copy of slide 4 with it.

    **The copy gets a fresh slide id**, never the source's and never one freed
    by an earlier delete, so anything holding a slide id across this call gets
    a clean miss rather than the wrong slide. SPEC §5.6.

    **Repeatable, not idempotent.** Calling it twice yields two distinct
    slides, and nobody should try to make it a no-op. What is idempotent is
    *allocation*: part names, relationship ids and slide ids are pure
    functions of package state plus the per-operation reservation set, so the
    second call cannot collide with the first.

    ``with_notes`` defaults to True because PowerPoint's Duplicate Slide keeps
    speaker notes, and because the costs are asymmetric: losing notes is
    unrecoverable, gaining them is one ``del`` away.

    Example:
        >>> from pptx import Presentation
        >>> prs = Presentation()
        >>> for _ in range(3):
        ...     _ = prs.slides.add_slide(prs.slide_layouts[6])
        >>> copy = duplicate_slide(prs, 0)
        >>> prs.slides.index(copy)
        1
    """
    index, slide, _sld_id = resolve_slide(prs, slide_or_index)

    lst = sld_id_lst(prs)
    count = len(lst) + 1  # the deck as it will be, with the copy in it
    target = index + 1 if to_index is None else (to_index + count if to_index < 0 else to_index)
    if not 0 <= target < count:
        raise SlideIndexError(
            f"to_index {to_index} is out of range for a deck of {count} slides "
            f"after the duplicate; valid positions are 0..{count - 1} "
            f"(or -{count}..-1)"
        )

    result = clone_part_graph(
        slide.part,
        into=prs.part.package,
        policy=ClonePolicy(with_notes=with_notes),
    )

    # `add_sldId` appends and allocates `max(used) + 1`, so the copy's id is
    # fresh by construction and the source keeps its own. It goes on the end
    # first and then moves, because allocation and placement are separate
    # concerns and only the allocator may choose the id.
    new_sld_id = lst.add_sldId(prs.part.relate_to(result.root, RT.SLIDE))
    if target != count - 1:
        lst.remove(new_sld_id)
        lst.insert(target, new_sld_id)

    insert_slide(prs, slide_id=new_sld_id.id, to_index=target)

    return cast("Slide", prs.slides[target])


__all__ = ["duplicate_slide"]

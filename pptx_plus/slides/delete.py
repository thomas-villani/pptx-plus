"""Remove a slide from its deck.

python-pptx's oldest open feature request, and the recipe in circulation is
wrong. ``prs.slides._sldIdLst.remove(sldId)`` unlinks the entry and stops:
the presentation part's relationship survives, so the slide part is still
reachable, still written to the saved package, and still named by any section
or custom show that contained it. The deck opens with a slide that has
vanished from the running order while continuing to inflate the file, and
PowerPoint reports the dangling section entry as damage.

What deletion actually is, in three steps:

1. Scrub the slide from ``p14:sectionLst`` and ``p:custShowLst``, while both
   of the identifiers naming it are still resolvable.
2. Remove its ``<p:sldId>`` from the running order.
3. Drop the presentation part's relationship to it.

Step 3 is what collects the slide. ``OpcPackage.save`` serializes a walk of
the relationship graph, so an unreferenced part is simply never written —
along with its notes slide, and any image reachable only through it. There is
no collector here and there should not be one (SPEC §3.5).

ECMA-376 Part 1 §19.2.1.34 (``p:sldIdLst``), §13.2 (part reachability).

SPEC §5.2. This module imports only from ``pptx_plus.core`` (SPEC §9.1).
"""

from __future__ import annotations

from typing import TYPE_CHECKING

from pptx_plus.core.oxml import remove
from pptx_plus.core.parts import drop_relationship
from pptx_plus.core.sections import scrub_slide
from pptx_plus.slides.resolve import resolve_slide

if TYPE_CHECKING:
    from pptx.presentation import Presentation
    from pptx.slide import Slide


def delete_slide(prs: Presentation, slide_or_index: Slide | int) -> None:
    """Delete a slide from the deck.

    Args:
        prs: The presentation to delete from.
        slide_or_index: The slide to delete, or its index. Negative indices
            count from the end.

    Raises:
        SlideIndexError: The index is outside the deck.
        SlideNotFoundError: The slide given is not in this deck.
        TypeError: ``slide_or_index`` is neither a ``Slide`` nor an ``int``.

    The slide, its notes slide, and every part reachable only through it stop
    being written on the next save. Parts still reachable from elsewhere —
    an image another slide also uses, the layout, the masters — are untouched.
    Neither behaviour is implemented here: both fall out of the package being
    serialized by a reachability walk (SPEC §3.5).

    **Deleting the last slide is allowed.** An empty ``<p:sldIdLst/>`` is
    schema-valid and PowerPoint opens a zero-slide deck, so this is a legal
    edge case rather than one to guard against.

    A second delete of the same slide raises ``SlideNotFoundError``. Because
    that subclasses ``KeyError``, callers who want idempotence can opt in with
    ``contextlib.suppress(KeyError)``; making the verb silently no-op instead
    would hide caller bugs from everyone to spare that one line. Note that
    ``delete_slide(prs, 0)`` twice deletes *two* slides — that is index
    semantics, not a contradiction. SPEC §5.6.

    The deleted ``Slide`` object stays alive and readable. Deletion detaches
    its part from the relationship graph; it destroys nothing, and this
    library will not poison an object it did not create. Use
    :func:`~pptx_plus.slides.resolve.contains` to ask whether a slide is still
    in the deck. SPEC §5.7.

    Note:
        Slide part names are **not** renumbered, so a gap is normal and
        harmless: content types are regenerated per part at save, and
        relationship targets are computed relative to the referring part at
        write time. SPEC §5.8.

    Example:
        >>> from pptx import Presentation
        >>> prs = Presentation()
        >>> for _ in range(3):
        ...     _ = prs.slides.add_slide(prs.slide_layouts[6])
        >>> delete_slide(prs, 1)
        >>> len(prs.slides)
        2
    """
    _index, _slide, sld_id = resolve_slide(prs, slide_or_index)
    r_id = sld_id.rId

    # Before anything is unlinked: the side-indexes name this slide by two
    # different identifiers and both have to still resolve. SPEC §3.3.
    scrub_slide(prs, slide_id=sld_id.id, r_id=r_id)

    remove(sld_id)

    # NOT `XmlPart.drop_rel`. That method is conditional -- it refuses when the
    # id is referenced twice or more in the part's XML -- and a slide that also
    # appeared in a custom show is referenced exactly twice. It would silently
    # do nothing and leave the slide serialized but unreachable from the deck.
    # (The scrub above removes the second reference, so `drop_rel` would in
    # fact work here today. It is still the wrong call: correctness would then
    # depend on the ordering of two statements rather than on the operation
    # doing what it says.) SPEC §4.6.
    drop_relationship(prs.part, r_id)


__all__ = ["delete_slide"]

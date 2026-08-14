"""Reorder a slide within its deck.

python-pptx has no API for this. ``Slides`` is a read-only sequence over
``p:sldIdLst``, and the element that determines running order is reachable
only as ``prs.slides._sldIdLst``.

Reordering is the one slide verb that touches no parts and no relationships:
a ``<p:sldId>`` moves within its parent and nothing else changes. What it does
have to maintain is ``p14:sectionLst`` — sections partition the running order
into contiguous runs, so a slide that moves across a boundary changes section.

ECMA-376 Part 1 §19.2.1.34 (``p:sldIdLst``), §19.2.1.33 (``p:sldId``).

SPEC §5.3. This module imports only from ``pptx_plus.core`` (SPEC §9.1).
"""

from __future__ import annotations

from typing import TYPE_CHECKING

from pptx_plus.core.oxml import sld_id_lst
from pptx_plus.core.sections import reorder_slide
from pptx_plus.slides.resolve import SlideIndexError, resolve_slide

if TYPE_CHECKING:
    from pptx.presentation import Presentation
    from pptx.slide import Slide


def move_slide(prs: Presentation, slide_or_index: Slide | int, to_index: int) -> None:
    """Move a slide to a new position in the deck.

    Args:
        prs: The presentation to reorder.
        slide_or_index: The slide to move, or its current index. Negative
            indices count from the end.
        to_index: The position the slide will occupy **in the resulting
            deck**. Negative values count from the end of the resulting deck.

    Raises:
        SlideIndexError: Either index is outside the deck.
        SlideNotFoundError: The slide given is not in this deck.
        TypeError: ``slide_or_index`` is neither a ``Slide`` nor an ``int``.

    ``to_index`` is the position in the deck *after* the move — ``list.insert``
    semantics against the list with the slide already removed:

    ========= ========================= ===============
    Deck      Call                      Result
    ========= ========================= ===============
    A B C D   ``move_slide(prs, 0, 2)`` B C **A** D
    A B C D   ``move_slide(prs, 2, 2)`` A B C D
    A B C D   ``move_slide(prs, 3, 0)`` **D** A B C
    ========= ========================= ===============

    The rejected reading — "index in the original list" — makes
    ``move_slide(prs, i, i + 1)`` a no-op, which surprises everyone.
    ``move_slide(prs, i, i)`` is a no-op for every valid ``i``, and that is the
    property that pins the semantics down. SPEC §5.5.

    Out-of-range raises rather than clamping. ``list.insert`` clamps; that is
    the right call for a list and the wrong one here, because clamping turns
    an off-by-one in a caller's loop into a silently misordered deck — the
    exact failure this library exists to make impossible.

    Example:
        >>> from pptx import Presentation
        >>> prs = Presentation()
        >>> for _ in range(3):
        ...     _ = prs.slides.add_slide(prs.slide_layouts[6])
        >>> first = prs.slides[0]
        >>> move_slide(prs, 0, 2)
        >>> prs.slides.index(first)
        2
    """
    index, _slide, sld_id = resolve_slide(prs, slide_or_index)
    lst = sld_id_lst(prs)
    count = len(lst)

    target = to_index + count if to_index < 0 else to_index
    if not 0 <= target < count:
        raise SlideIndexError(
            f"to_index {to_index} is out of range for a deck of {count} "
            f"slide{'' if count == 1 else 's'}; valid positions are "
            f"0..{count - 1} (or -{count}..-1)"
        )
    if target == index:
        return

    # Remove first, then insert: `to_index` is a position in the resulting
    # deck, so the arithmetic is against the list without this slide in it.
    lst.remove(sld_id)
    lst.insert(target, sld_id)

    # `.id` rather than `.get("id")`: the typed accessor, and the one the
    # upstream-surface guard covers. Sections key on this value, not on the
    # relationship id sitting next to it on the same element (SPEC §3.3).
    reorder_slide(prs, slide_id=sld_id.id, to_index=target)


__all__ = ["move_slide"]

"""Argument normalization for the slide verbs.

python-pptx's ``Slides`` collection indexes and iterates, but offers no way to
go from a ``Slide`` back to its position, and no way at all to reach the
``<p:sldId>`` element that actually holds the slide's identity. Every verb in
this package needs all three — the index for ordering arithmetic, the ``Slide``
to return or report, and the ``<p:sldId>`` to move, remove, or read the slide
id from.

Resolving them together is not a convenience. Resolved separately they mean
three walks of ``sldIdLst`` and three places to get index arithmetic wrong,
and the failure mode of getting it wrong is operating on the neighbouring
slide — which looks like a working program.

ECMA-376 Part 1 §19.2.1.34 (``p:sldIdLst``).

SPEC §5.1. This module imports only from ``pptx_plus.core`` (SPEC §9.1).
"""

from __future__ import annotations

from typing import TYPE_CHECKING, cast

# Imported at runtime, not under TYPE_CHECKING, because `resolve_slide`
# needs it for an `isinstance` check -- the argument may be anything a
# caller passed, and the error message names the type received.
from pptx.slide import Slide

from pptx_plus.core.errors import PptxPlusError
from pptx_plus.core.oxml import sld_id_lst

if TYPE_CHECKING:
    from pptx.oxml.presentation import CT_SlideId
    from pptx.presentation import Presentation


class SlideNotFoundError(PptxPlusError, KeyError):
    """Raised when a ``Slide`` is not part of the presentation given.

    Subclasses ``KeyError`` both for ``except`` ergonomics and because it is
    what gives callers opt-in idempotence:
    ``with contextlib.suppress(KeyError): delete_slide(prs, slide)``. Making
    the verbs silently no-op instead would hide caller bugs from everyone to
    spare that one line. SPEC §5.6.
    """


class SlideIndexError(PptxPlusError, IndexError):
    """Raised for a slide index outside the deck.

    Subclasses ``IndexError`` per SPEC §9.7, so code that already handles
    out-of-range indexing keeps working.
    """


def resolve_slide(prs: Presentation, slide_or_index: Slide | int) -> tuple[int, Slide, CT_SlideId]:
    """Normalize a slide argument to its index, object, and ``<p:sldId>``.

    Args:
        prs: The presentation to resolve against.
        slide_or_index: A :class:`~pptx.slide.Slide` belonging to ``prs``, or
            an index into it. Negative indices count from the end, following
            list semantics.

    Returns:
        ``(index, slide, sld_id)`` — the position in the deck, the slide
        itself, and the ``<p:sldId>`` element that holds its deck-scoped id and
        its relationship id.

    Raises:
        SlideIndexError: The index is outside the deck.
        SlideNotFoundError: The slide is not in this deck.
        TypeError: The argument is neither a ``Slide`` nor an ``int``.

    A ``Slide`` is located by the identity of its underlying ``<p:sld>``
    element rather than by equality. python-pptx defines no ``__eq__`` on
    ``Slide``, so equality would be identity anyway — but going through the
    element says what is actually meant and survives a future upstream change
    that starts caching wrapper objects differently.
    """
    lst = sld_id_lst(prs)
    count = len(lst)

    if isinstance(slide_or_index, int):
        index = slide_or_index + count if slide_or_index < 0 else slide_or_index
        if not 0 <= index < count:
            raise SlideIndexError(
                f"slide index {slide_or_index} is out of range for a deck of {count} "
                f"slide{'' if count == 1 else 's'}"
            )
        return index, prs.slides[index], cast("CT_SlideId", lst[index])

    if not isinstance(slide_or_index, Slide):
        raise TypeError(f"expected a Slide or an int, got {type(slide_or_index).__name__}")

    element = slide_or_index.element
    for index in range(count):
        candidate = prs.slides[index]
        if candidate.element is element:
            return index, candidate, cast("CT_SlideId", lst[index])
    raise SlideNotFoundError(
        "the slide given is not in this presentation. A slide deleted earlier "
        "stays alive and readable but is no longer part of the deck (SPEC §5.7)."
    )


def slide_index(prs: Presentation, slide_or_index: Slide | int) -> int:
    """Return a slide's position in the deck.

    Args:
        prs: The presentation.
        slide_or_index: A slide belonging to ``prs``, or an index into it.

    Returns:
        The zero-based index, with a negative input normalized.

    Raises:
        SlideIndexError: The index is outside the deck.
        SlideNotFoundError: The slide is not in this deck.
    """
    return resolve_slide(prs, slide_or_index)[0]


def contains(prs: Presentation, slide: Slide) -> bool:
    """Return whether a slide is still part of the deck.

    Args:
        prs: The presentation.
        slide: The slide to look for.

    Returns:
        True if the slide is in ``prs``.

    The non-raising counterpart to :func:`slide_index`, for the check that is
    a question rather than an error. A deleted ``Slide`` object remains alive
    and readable — deletion detaches its part from the relationship graph and
    destroys nothing (SPEC §5.7) — so "do I still have this slide?" has no
    answer that can be read off the object itself.
    """
    element = slide.element
    return any(prs.slides[index].element is element for index in range(len(sld_id_lst(prs))))


__all__ = [
    "SlideIndexError",
    "SlideNotFoundError",
    "contains",
    "resolve_slide",
    "slide_index",
]

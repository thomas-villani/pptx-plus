"""Grading the harness itself -- SPEC §10.5.

The assertions in `pptx_plus._testing` are the oracle every later phase is
graded against, so they are built before the code they grade and graded first
themselves. Three steps, in order:

1. **They pass on valid input.** Every generated fixture, unmodified, must
   satisfy the whole battery. An oracle that fires on a correct package is
   worse than no oracle -- it trains you to ignore it.
2. **They fail on the known-broken recipes.** Each naive recipe from SPEC §3.6
   is performed here deliberately, and the specific invariant it violates is
   asserted to fire. A harness that cannot detect the known-broken recipe
   cannot be trusted to grade the correct one.
3. **Each assertion is individually falsifiable.** A targeted corruption per
   invariant, so no assertion is passing vacuously.
"""

from __future__ import annotations

import io
import zipfile
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

from pptx_plus._testing import (
    SavedPackage,
    assert_all_xml_parses,
    assert_content_types_complete,
    assert_in_package,
    assert_no_unclaimed_rid_literals,
    assert_not_in_package,
    assert_package_integrity,
    assert_partnames_unique,
    assert_parts_disjoint,
    assert_parts_shared,
    assert_rel_ids_resolve,
    assert_sections_consistent,
    assert_slide_ids_valid,
    assert_slide_rels_consistent,
    roundtrip,
    saved,
)
from pptx_plus.core.oxml import part_root, sld_id_lst
from pptx_plus.core.parts import clone_part
from tests.fixtures.build_decks import BUILDERS

ALL_FIXTURES = sorted(BUILDERS)


# ---------------------------------------------------------------------------
# Step 1 -- the battery passes on valid input
# ---------------------------------------------------------------------------


@pytest.mark.parametrize("name", ALL_FIXTURES)
def test_every_fixture_passes_the_battery(deck, name: str) -> None:
    """Unmodified, every generated deck is a valid package."""
    assert_package_integrity(SavedPackage(Path(deck(name)).read_bytes()))


@pytest.mark.parametrize("name", ALL_FIXTURES)
def test_every_fixture_survives_an_untouched_round_trip(deck, name: str) -> None:
    """Opening and re-saving without edits still yields a valid package.

    python-pptx re-serializes every part it parsed, so this proves the battery
    tolerates its output as well as PowerPoint's -- a distinction that matters,
    because every later test asserts against a re-saved deck.
    """
    assert_package_integrity(saved(Presentation(str(deck(name)))))


def test_the_gap_rids_deck_really_has_a_gap(deck) -> None:
    """The fixture is only useful if the gap is actually there."""
    pkg = SavedPackage(Path(deck("gap_rids")).read_bytes())
    assert sorted(pkg.rels(pkg.slide_partnames[0])) == ["rId1", "rId3"]


def test_a_relationship_id_gap_is_legal(deck) -> None:
    """A gap is unusual, not invalid -- the battery must not object to it."""
    assert_package_integrity(SavedPackage(Path(deck("gap_rids")).read_bytes()))


def test_an_empty_rel_id_is_not_a_dangling_reference(deck) -> None:
    """`<a:hlinkClick r:id="">` is ordinary, and must not read as dangling."""
    pkg = SavedPackage(Path(deck("hyperlink")).read_bytes())
    blob = pkg.blob_of(pkg.slide_partnames[0])
    assert b'r:id=""' in blob
    assert_rel_ids_resolve(pkg)


def test_an_external_relationship_resolves(deck) -> None:
    """An external target has no part behind it but still has an id."""
    pkg = SavedPackage(Path(deck("hyperlink")).read_bytes())
    external = [rel for rel in pkg.rels(pkg.slide_partnames[0]).values() if rel.is_external]
    assert [rel.partname for rel in external] == [None]


def test_slide_order_comes_from_the_slide_id_list(deck) -> None:
    """Order is `sldIdLst`'s to define; the part name says nothing about it."""
    pkg = SavedPackage(Path(deck("simple")).read_bytes())
    assert pkg.slide_partnames == [
        "/ppt/slides/slide1.xml",
        "/ppt/slides/slide2.xml",
        "/ppt/slides/slide3.xml",
    ]


def test_an_empty_deck_passes(deck) -> None:
    """Zero slides is schema-valid, so the battery must accept it."""
    assert_package_integrity(SavedPackage(Path(deck("empty")).read_bytes()))


def test_roundtrip_returns_a_usable_presentation(deck) -> None:
    assert len(roundtrip(Presentation(str(deck("simple")))).slides) == 3


# ---------------------------------------------------------------------------
# Step 2 -- the battery fails on the known-broken recipes (SPEC §3.6)
# ---------------------------------------------------------------------------


def _naive_delete(prs: Presentation, index: int) -> None:
    """The recipe in circulation: unlink the sldId and stop.

    It leaves the presentation part's relationship in place, so the slide part
    is still reachable and is still written to the package.
    """
    lst = sld_id_lst(prs)
    lst.remove(lst[index])


def _naive_duplicate(prs: Presentation, index: int) -> None:
    """Clone the slide part and attach it, skipping the relationship rewrite.

    This is `deepcopy(slide.element)` in effect: the copied XML keeps the
    source's relationship ids while the new part has relationships of its own
    -- here, none at all. Every `r:embed` in it is now dangling.
    """
    source = prs.slides[index].part
    clone = clone_part(source, into=source.package, reserved=set())
    sld_id_lst(prs).add_sldId(prs.part.relate_to(clone, RT.SLIDE))


def test_naive_delete_is_caught(deck) -> None:
    """The delete recipe leaves an orphaned slide relationship -- I2."""
    prs = Presentation(str(deck("simple")))
    _naive_delete(prs, 1)
    with pytest.raises(AssertionError, match="I2"):
        assert_package_integrity(saved(prs))


def test_naive_delete_leaves_the_slide_part_in_the_package(deck) -> None:
    """Naming the mechanism, not just the symptom: the part is still written."""
    prs = Presentation(str(deck("simple")))
    _naive_delete(prs, 1)
    assert_in_package(saved(prs), "/ppt/slides/slide2.xml")


def test_naive_duplicate_is_caught(deck) -> None:
    """The duplicate recipe leaves dangling relationship references -- I4."""
    prs = Presentation(str(deck("picture")))
    _naive_duplicate(prs, 0)
    with pytest.raises(AssertionError, match="I4"):
        assert_package_integrity(saved(prs))


def test_naive_duplicate_names_the_dangling_attribute(deck) -> None:
    """The failure has to say *what* dangled, or it cannot be acted on."""
    prs = Presentation(str(deck("picture")))
    _naive_duplicate(prs, 0)
    with pytest.raises(AssertionError, match="dangling relationship reference"):
        assert_rel_ids_resolve(saved(prs))


def test_naive_delete_of_a_sectioned_deck_is_caught(deck) -> None:
    """A stale section entry is the corruption no generated fixture shows -- I6.

    Sections are keyed on slide id and unmodelled by python-pptx, so this is
    the failure that passes every other assertion in the battery and still
    produces a repair prompt on a real user's deck.
    """
    prs = Presentation(str(deck("sections")))
    _naive_delete(prs, 1)
    with pytest.raises(AssertionError, match="I6"):
        assert_sections_consistent(saved(prs))


def test_naive_delete_of_a_custom_show_is_caught(deck) -> None:
    """A custom show keys on `r:id`, so it fails differently -- also I6."""
    prs = Presentation(str(deck("custom_show")))
    lst = sld_id_lst(prs)
    # Slide 0 is *in* the show; deleting a slide the show never named would
    # not break it, and the test would pass without proving anything.
    prs.part.rels.pop(lst[0].rId)
    lst.remove(lst[0])
    with pytest.raises(AssertionError, match="I6"):
        assert_sections_consistent(saved(prs))


# ---------------------------------------------------------------------------
# Step 3 -- each assertion is individually falsifiable
# ---------------------------------------------------------------------------


def test_slide_ids_out_of_range_are_caught(deck) -> None:
    prs = Presentation(str(deck("simple")))
    sld_id_lst(prs)[0].set("id", "12")
    with pytest.raises(AssertionError, match=r"outside the legal range"):
        assert_slide_ids_valid(saved(prs))


def test_duplicate_slide_ids_are_caught(deck) -> None:
    prs = Presentation(str(deck("simple")))
    lst = sld_id_lst(prs)
    lst[1].set("id", lst[0].get("id"))
    with pytest.raises(AssertionError, match="duplicate p:sldId/@id"):
        assert_slide_ids_valid(saved(prs))


def test_a_slide_id_naming_no_relationship_is_caught(deck) -> None:
    prs = Presentation(str(deck("simple")))
    sld_id_lst(prs)[0].rId = "rId999"
    with pytest.raises(AssertionError, match="I1"):
        assert_slide_rels_consistent(saved(prs))


def test_a_slide_id_pointing_at_a_non_slide_is_caught(deck) -> None:
    """`r:id` resolving to *something* is not enough; it must be a slide."""
    prs = Presentation(str(deck("simple")))
    master_rid = next(r_id for r_id, rel in prs.part.rels.items() if rel.reltype == RT.SLIDE_MASTER)
    sld_id_lst(prs)[0].rId = master_rid
    with pytest.raises(AssertionError, match="not a slide"):
        assert_slide_rels_consistent(saved(prs))


def test_an_unclaimed_rid_literal_is_caught(deck) -> None:
    """An attribute holding an id the rewriter would not know to rewrite."""
    prs = Presentation(str(deck("simple")))
    part_root(prs.slides[0].part).set("name", "rId7")
    with pytest.raises(AssertionError, match="unclaimed relationship-id literal"):
        assert_no_unclaimed_rid_literals(saved(prs))


def test_a_qualified_rel_id_is_not_reported_as_unclaimed(deck) -> None:
    """The detector must not fire on the ids the sweep already handles."""
    assert_no_unclaimed_rid_literals(SavedPackage(Path(deck("picture")).read_bytes()))


@pytest.mark.filterwarnings("ignore:Duplicate name:UserWarning")
def test_a_duplicate_partname_is_caught(deck) -> None:
    """`zipfile` warns as it writes the duplicate; that is the point of the test."""
    pkg = _corrupt_zip(deck("simple"), duplicate_entry="ppt/slides/slide1.xml")
    with pytest.raises(AssertionError, match="I5: duplicate part name"):
        assert_partnames_unique(pkg)


def test_a_missing_content_type_is_caught(deck) -> None:
    """A part written under an extension nobody declared.

    Not a dropped slide Override -- `Default Extension="xml"` catches that and
    the package stays technically complete. This is the case that genuinely
    occurs when a part type is added without registering it.
    """
    pkg = _corrupt_zip(deck("simple"), add_entry="ppt/media/thing.qqq")
    with pytest.raises(AssertionError, match="no content type declared"):
        assert_content_types_complete(pkg)


def test_malformed_xml_is_caught(deck) -> None:
    pkg = _corrupt_zip(deck("simple"), truncate="ppt/slides/slide1.xml")
    with pytest.raises(AssertionError, match="not well-formed XML"):
        assert_all_xml_parses(pkg)


def test_assert_not_in_package_is_falsifiable(deck) -> None:
    pkg = SavedPackage(Path(deck("simple")).read_bytes())
    with pytest.raises(AssertionError, match="still in the saved package"):
        assert_not_in_package(pkg, "/ppt/slides/slide1.xml")


def test_assert_in_package_is_falsifiable(deck) -> None:
    pkg = SavedPackage(Path(deck("simple")).read_bytes())
    with pytest.raises(AssertionError, match="not in the saved package"):
        assert_in_package(pkg, "/ppt/slides/slide99.xml")


# ---------------------------------------------------------------------------
# I7 -- sharing and ownership, on a deck no verb has touched yet
# ---------------------------------------------------------------------------


def test_two_slides_showing_one_image_are_shared(deck) -> None:
    """python-pptx dedupes on image bytes, so this holds before any verb runs."""
    pkg = SavedPackage(Path(deck("shared_picture")).read_bytes())
    first, second = pkg.slide_partnames
    assert_parts_shared(pkg, first, second, reltypes=[RT.IMAGE])


def test_assert_parts_shared_is_falsifiable(deck) -> None:
    pkg = SavedPackage(Path(deck("simple")).read_bytes())
    first, second = pkg.slide_partnames[:2]
    with pytest.raises(AssertionError, match="no relationships of the given types"):
        assert_parts_shared(pkg, first, second, reltypes=[RT.IMAGE])


def test_two_charts_own_separate_workbooks(deck) -> None:
    """Each chart part owns its own embedded workbook -- the I7 ownership case."""
    pkg = SavedPackage(Path(deck("two_charts")).read_bytes())
    charts = sorted(_related(pkg, pkg.slide_partnames[0], RT.CHART))
    assert len(charts) == 2
    assert_parts_disjoint(pkg, charts[0], charts[1], reltypes=[RT.PACKAGE])


def test_assert_parts_disjoint_is_falsifiable(deck) -> None:
    """A part is never disjoint from itself."""
    pkg = SavedPackage(Path(deck("shared_picture")).read_bytes())
    first, second = pkg.slide_partnames
    with pytest.raises(AssertionError, match="share owned part"):
        assert_parts_disjoint(pkg, first, second, reltypes=[RT.IMAGE])


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _related(pkg: SavedPackage, partname: str, reltype: str) -> set[str]:
    return {rel.partname for rel in pkg.targets(partname, [reltype]) if rel.partname}


def _corrupt_zip(
    source: Path,
    *,
    duplicate_entry: str | None = None,
    add_entry: str | None = None,
    truncate: str | None = None,
) -> SavedPackage:
    """Rebuild a package with one targeted defect, for falsifying an assertion.

    Zip-level rather than model-level because these are defects python-pptx
    will not produce and cannot represent -- which is the point: the harness
    reads packages it did not write.
    """
    with zipfile.ZipFile(source) as src:
        entries = [(item.filename, src.read(item.filename)) for item in src.infolist()]

    out = io.BytesIO()
    with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as dst:
        for name, data in entries:
            if name == truncate:
                data = data[: len(data) // 2]
            dst.writestr(name, data)
            if name == duplicate_entry:
                dst.writestr(name, data)
        if add_entry:
            dst.writestr(add_entry, b"undeclared")
    return SavedPackage(out.getvalue())

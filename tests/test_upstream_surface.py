"""SPEC §14.2 — every python-pptx attribute pptx_plus depends on still exists.

One test per attribute, so bumping the python-pptx pin produces a red test
naming the thing that moved rather than a mysterious failure inside a clone.

The tests below the parametrized check are the interesting ones: they pin
*behaviours* this library's correctness rests on, not merely the existence of
a name. An attribute that still exists but has changed semantics is the
failure mode a presence check cannot see.
"""

from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.opc.package import Part, XmlPart, _Relationships

from pptx_plus.core._compat import (
    REQUIRED_SURFACE,
    UpstreamSurfaceError,
    _resolve,
    check_upstream_surface,
)


@pytest.mark.parametrize(
    ("module_path", "dotted", "why"),
    REQUIRED_SURFACE,
    ids=[f"{m.rsplit('.', 1)[-1]}.{d}" for m, d, _ in REQUIRED_SURFACE],
)
def test_required_attribute_exists(module_path: str, dotted: str, why: str) -> None:
    assert _resolve(module_path, dotted), f"python-pptx no longer provides {dotted} — {why}"


def test_check_passes_on_the_pinned_version() -> None:
    check_upstream_surface()  # must not raise


def test_check_reports_a_missing_attribute(monkeypatch: pytest.MonkeyPatch) -> None:
    """Guard the guard — the check must be able to fail."""
    monkeypatch.setattr(
        "pptx_plus.core._compat.REQUIRED_SURFACE",
        (("pptx.opc.package", "Part.no_such_attribute", "a fabricated requirement"),),
    )
    with pytest.raises(UpstreamSurfaceError, match="no_such_attribute"):
        check_upstream_surface()


def test_upstream_surface_error_is_a_runtime_error() -> None:
    assert issubclass(UpstreamSurfaceError, RuntimeError)


# --------------------------------------------------------------------------
# Behaviours, not just names. These are the assumptions the design rests on.
# --------------------------------------------------------------------------


def test_no_part_subclass_overrides_load() -> None:
    """The clone primitive's whole basis — SPEC §4.6.

    `type(src).load(partname, content_type, package, blob)` clones a slide, a
    chart, an image, and a blob-only SmartArt part with one expression *only*
    because every part class inherits `load` unchanged. If a subclass ever
    overrides it, the clone path needs to know rather than silently taking a
    different code path for one part type.
    """
    import importlib
    import pkgutil

    import pptx

    overrides: list[str] = []
    for module_info in pkgutil.walk_packages(pptx.__path__, "pptx."):
        try:
            module = importlib.import_module(module_info.name)
        except ImportError:  # pragma: no cover - defensive
            continue
        overrides.extend(
            f"{module_info.name}.{name}"
            for name, obj in vars(module).items()
            if isinstance(obj, type)
            and issubclass(obj, Part)
            and obj not in (Part, XmlPart)
            and "load" in vars(obj)
        )
    assert not overrides, f"a Part subclass now overrides load(): {sorted(set(overrides))}"


def test_next_rId_fills_gaps() -> None:
    """Why the relationship-id remap is mandatory — SPEC §4.4.

    A slide whose rels are {rId1, rId3} — the ordinary result of deleting a
    shape in PowerPoint — clones to {rId1, rId2}, so `r:embed="rId3"` in the
    copied XML would dangle. This is the *most common* trigger in real decks,
    and the simple case that hides it produces an identity map.
    """
    prs = Presentation()
    part = prs.slides.add_slide(prs.slide_layouts[6]).part
    rels = part.rels
    while len(rels) < 3:
        rels.get_or_add_ext_rel("http://example.invalid/t", f"http://example.invalid/{len(rels)}")
    second = sorted(rels)[1]
    rels.pop(second)
    assert rels._next_rId == second


def test_get_or_add_dedupes_on_reltype_and_target() -> None:
    """Why RelMap is many-to-one — SPEC §4.4.

    Two source relationships with the same type and target collapse to one on
    the clone, so two old ids map to one new id. Code asserting the map is a
    bijection would fail on a valid deck.
    """
    prs = Presentation()
    part = prs.slides.add_slide(prs.slide_layouts[6]).part
    first = part.rels.get_or_add_ext_rel("http://example.invalid/t", "http://example.invalid/x")
    second = part.rels.get_or_add_ext_rel("http://example.invalid/t", "http://example.invalid/x")
    assert first == second


def test_drop_rel_is_conditional() -> None:
    """Why delete_slide uses rels.pop() instead — SPEC §5.2.

    `XmlPart.drop_rel` refuses to drop a relationship referenced twice or more
    in the part's XML, counting over `//@r:id`. For a slide that also appears
    in a custom show, that count is 2 and the call silently does nothing.
    """
    assert "_rel_ref_count" in dir(XmlPart)
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    pres_part = prs.part
    rid = prs.slides._sldIdLst[0].rId
    assert pres_part._rel_ref_count(rid) == 1
    pres_part.drop_rel(rid)
    assert rid not in pres_part.rels


def test_rels_is_a_public_mapping() -> None:
    """`Part.rels` is documented upstream as necessarily public."""
    assert isinstance(Presentation().part.rels, _Relationships)

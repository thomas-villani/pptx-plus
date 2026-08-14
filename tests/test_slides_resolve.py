"""Argument normalization for the slide verbs -- SPEC §5.1."""

from __future__ import annotations

import pytest
from pptx import Presentation

from pptx_plus.core.oxml import sld_id_lst
from pptx_plus.slides.resolve import (
    SlideIndexError,
    SlideNotFoundError,
    contains,
    resolve_slide,
    slide_index,
)


@pytest.fixture
def prs(deck) -> Presentation:
    return Presentation(str(deck("simple")))


# ---------------------------------------------------------------------------
# By index
# ---------------------------------------------------------------------------


def test_an_index_resolves_to_that_position(prs: Presentation) -> None:
    assert resolve_slide(prs, 1)[0] == 1


def test_an_index_resolves_to_the_matching_slide(prs: Presentation) -> None:
    assert resolve_slide(prs, 1)[1] is prs.slides[1]


def test_an_index_resolves_to_the_matching_sld_id(prs: Presentation) -> None:
    assert resolve_slide(prs, 1)[2] is sld_id_lst(prs)[1]


@pytest.mark.parametrize(("given", "expected"), [(-1, 2), (-2, 1), (-3, 0)])
def test_a_negative_index_counts_from_the_end(prs: Presentation, given: int, expected: int) -> None:
    assert resolve_slide(prs, given)[0] == expected


@pytest.mark.parametrize("given", [3, 99, -4])
def test_an_out_of_range_index_raises(prs: Presentation, given: int) -> None:
    with pytest.raises(SlideIndexError, match="out of range"):
        resolve_slide(prs, given)


def test_the_range_error_reports_the_deck_size(prs: Presentation) -> None:
    """The message has to say what the valid range *was*, or it cannot be acted on."""
    with pytest.raises(SlideIndexError, match="deck of 3 slides"):
        resolve_slide(prs, 7)


def test_a_one_slide_deck_says_slide_not_slides() -> None:
    """Pluralization, because an error message is read by a person."""
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[6])
    with pytest.raises(SlideIndexError, match="deck of 1 slide$"):
        resolve_slide(presentation, 5)


def test_an_empty_deck_rejects_every_index() -> None:
    with pytest.raises(SlideIndexError):
        resolve_slide(Presentation(), 0)


# ---------------------------------------------------------------------------
# By slide
# ---------------------------------------------------------------------------


def test_a_slide_resolves_to_its_position(prs: Presentation) -> None:
    assert resolve_slide(prs, prs.slides[2])[0] == 2


def test_a_slide_is_matched_by_element_identity(prs: Presentation) -> None:
    """Not by equality -- python-pptx defines no `__eq__` on `Slide`."""
    assert resolve_slide(prs, prs.slides[1])[2] is sld_id_lst(prs)[1]


def test_a_slide_from_another_deck_raises(prs: Presentation, deck) -> None:
    other = Presentation(str(deck("simple")))
    with pytest.raises(SlideNotFoundError, match="not in this presentation"):
        resolve_slide(prs, other.slides[0])


def test_slide_not_found_is_a_key_error() -> None:
    """What makes `contextlib.suppress(KeyError)` an opt-in idempotence. SPEC §5.6."""
    assert issubclass(SlideNotFoundError, KeyError)


def test_slide_index_error_is_an_index_error() -> None:
    assert issubclass(SlideIndexError, IndexError)


# ---------------------------------------------------------------------------
# Bad types
# ---------------------------------------------------------------------------


@pytest.mark.parametrize("given", ["0", None, 1.0, object()])
def test_an_unsupported_type_raises_type_error(prs: Presentation, given: object) -> None:
    with pytest.raises(TypeError, match="expected a Slide or an int"):
        resolve_slide(prs, given)  # type: ignore[arg-type]


def test_the_type_error_names_the_type_received(prs: Presentation) -> None:
    with pytest.raises(TypeError, match="got str"):
        resolve_slide(prs, "1")  # type: ignore[arg-type]


# ---------------------------------------------------------------------------
# slide_index and contains
# ---------------------------------------------------------------------------


def test_slide_index_returns_the_position(prs: Presentation) -> None:
    assert slide_index(prs, prs.slides[2]) == 2


def test_slide_index_normalizes_a_negative_index(prs: Presentation) -> None:
    assert slide_index(prs, -1) == 2


def test_contains_is_true_for_a_slide_in_the_deck(prs: Presentation) -> None:
    assert contains(prs, prs.slides[0]) is True


def test_contains_is_false_for_a_foreign_slide(prs: Presentation, deck) -> None:
    other = Presentation(str(deck("simple")))
    assert contains(prs, other.slides[0]) is False


def test_contains_does_not_raise_where_slide_index_would(prs: Presentation, deck) -> None:
    """The non-raising counterpart, for the check that is a question."""
    other = Presentation(str(deck("simple")))
    assert contains(prs, other.slides[0]) is False
    with pytest.raises(SlideNotFoundError):
        slide_index(prs, other.slides[0])

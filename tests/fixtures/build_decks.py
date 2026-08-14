"""Deterministic fixture decks, generated on demand and never committed.

Every deck here is authored by python-pptx, or by zip surgery on something
python-pptx authored. Nothing in this file writes into the source tree: the
builders take a destination directory, and `main()` defaults to a temp one.
A generated artifact checked into a repository goes stale silently and gets
believed anyway.

The decks python-pptx *cannot* author -- SmartArt, embedded video, sections,
custom shows -- are committed binaries under `pptx_samples/` instead, with
their provenance recorded there. SPEC §10.6.

Run standalone to eyeball the output:

    uv run python -m tests.fixtures.build_decks --out build/fixtures
"""

from __future__ import annotations

import argparse
import io
import re
import tempfile
import zipfile
from collections.abc import Callable
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.slide import Slide
from pptx.util import Inches, Pt

from pptx_plus.core.oxml import sub

#: A 1x1 transparent GIF -- the smallest thing python-pptx accepts as a
#: picture. Fixtures care that an image part exists and is shared or cloned,
#: never what it depicts, so the cheapest possible image is the right one.
GIF = (
    b"GIF89a\x01\x00\x01\x00\x80\x00\x00\x00\x00\x00\xff\xff\xff!"
    b"\xf9\x04\x01\x00\x00\x00\x00,\x00\x00\x00\x00\x01\x00\x01\x00\x00\x02\x02D\x01\x00;"
)

#: Layout 6 of the default template is "Blank" -- no placeholders, so a slide
#: built on it contains exactly what the fixture puts there.
BLANK = 6
#: Layout 5 is "Title Only", used where a fixture wants readable slide titles.
TITLE_ONLY = 5


def _titled(prs: Presentation, title: str) -> Slide:
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_ONLY])
    slide.shapes.title.text = title
    return slide


# ---------------------------------------------------------------------------
# Builders
# ---------------------------------------------------------------------------


def build_empty(path: Path) -> Path:
    """A deck with no slides at all.

    An empty `<p:sldIdLst/>` is schema-valid and PowerPoint opens a zero-slide
    deck, so this is the legal edge case that proves deleting the last slide
    needs no guard.
    """
    Presentation().save(path)
    return path


def build_simple(path: Path, count: int = 3) -> Path:
    """`count` titled slides, no media. The index-arithmetic workhorse."""
    prs = Presentation()
    for index in range(count):
        _titled(prs, f"Slide {index + 1}")
    prs.save(path)
    return path


def build_picture(path: Path) -> Path:
    """One slide bearing one image: the `a:blip/@r:embed` rewrite target."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
    slide.shapes.add_picture(io.BytesIO(GIF), Inches(1), Inches(1))
    prs.save(path)
    return path


def build_shared_picture(path: Path) -> Path:
    """Two slides showing the *same* image part.

    python-pptx dedupes image parts on their SHA1, so adding the same bytes
    twice yields one part with two referrers. That is the deck where deleting
    slide 1 must leave the image alone -- and where a naive "delete the slide's
    parts" implementation breaks slide 2.
    """
    prs = Presentation()
    for _ in range(2):
        slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
        slide.shapes.add_picture(io.BytesIO(GIF), Inches(1), Inches(1))
    prs.save(path)
    return path


def build_notes(path: Path) -> Path:
    """Three slides, the middle one carrying speaker notes.

    The notes slide holds an `RT.SLIDE` relationship *back* to its slide, so
    this is the deck that proves the clone map resolves a cycle rather than
    recursing forever.
    """
    prs = Presentation()
    for index in range(3):
        slide = _titled(prs, f"Slide {index + 1}")
        if index == 1:
            slide.notes_slide.notes_text_frame.text = "Speaker notes for slide 2."
    prs.save(path)
    return path


def build_chart(path: Path) -> Path:
    """One slide with one chart, which owns an embedded xlsx workbook.

    The chart part is an XML part with its own relationship to a blob part --
    a two-level sub-graph, and the case that proves cloning is recursive
    rather than one level deep.
    """
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
    data = CategoryChartData()
    data.categories = ["Q1", "Q2", "Q3"]
    data.add_series("Revenue", (10.0, 20.0, 30.0))
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1), Inches(6), Inches(4), data
    )
    prs.save(path)
    return path


def build_two_charts(path: Path) -> Path:
    """One slide with two charts: the part-name reservation regression.

    Without a reservation set, both cloned charts are allocated
    `/ppt/charts/chart3.xml` -- `next_partname` cannot see a part that has
    been constructed but not yet related to anything -- and the second zip
    entry silently overwrites the first. SPEC §4.6.
    """
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
    for column in range(2):
        data = CategoryChartData()
        data.categories = ["A", "B"]
        data.add_series(f"S{column}", (1.0, 2.0))
        slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(1 + 4 * column),
            Inches(1),
            Inches(3.5),
            Inches(3),
            data,
        )
    prs.save(path)
    return path


def build_hyperlink(path: Path) -> Path:
    """A slide with an external hyperlink and an action-only one.

    Two things worth pinning. The external relationship has `TargetMode=
    "External"` and no part behind it, so it must be re-minted rather than
    cloned. The action-only run carries `<a:hlinkClick r:id="">` -- an *empty*
    relationship id, which is ordinary in real decks and must not be mistaken
    for a dangling reference.

    The empty one is authored by hand because python-pptx has no API for it:
    `run.hyperlink.address` only ever mints a relationship. PowerPoint writes
    it constantly -- every "go to next slide" action button is one -- so a
    harness that has never seen it would fail on the first real deck.
    """
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    paragraph = box.text_frame.paragraphs[0]

    linked = paragraph.add_run()
    linked.text = "python-pptx"
    linked.font.size = Pt(18)
    linked.hyperlink.address = "https://github.com/scanny/python-pptx"

    action = paragraph.add_run()
    action.text = " (next slide)"
    action.font.size = Pt(18)
    sub(
        action._r.get_or_add_rPr(),  # noqa: SLF001 - no public path to an action-only link
        "a:hlinkClick",
        **{"r:id": "", "action": "ppaction://hlinkshowjump?jump=nextslide"},
    )

    prs.save(path)
    return path


# ---------------------------------------------------------------------------
# Zip surgery -- decks python-pptx will not author
# ---------------------------------------------------------------------------


def build_gap_rids(path: Path) -> Path:
    """A picture deck whose slide relationship ids have a gap: {rId1, rId3}.

    **The primary regression for the relationship remap.** In a freshly
    authored deck the cloned ids come out identical to the source's, so an
    implementation that skips the rewrite entirely still passes -- by
    accident. This deck removes the accident.

    `_next_rId` documents itself as filling gaps, so a source with `{rId1,
    rId3}` clones to `{rId1, rId2}` and an unrewritten `r:embed="rId3"`
    dangles. The gap is not exotic: it is what PowerPoint leaves behind every
    time a user deletes a shape.

    The deck is *legal* -- gaps in relationship ids are perfectly valid -- so
    it must pass the whole integrity battery unmodified.
    """
    source = build_picture(path.with_name(f"{path.stem}-source.pptx"))
    try:
        _rewrite_zip(source, path, _renumber_slide_image_rel)
    finally:
        source.unlink(missing_ok=True)
    return path


def _renumber_slide_image_rel(name: str, data: bytes) -> bytes:
    """Rename the slide's image relationship from rId2 to rId3, creating a gap.

    Deliberately narrow: it touches only slide1 and its `.rels`, and only the
    one id, so the resulting deck differs from its source in exactly the
    property under test.
    """
    if name == "ppt/slides/_rels/slide1.xml.rels":
        return re.sub(rb'Id="rId2"', b'Id="rId3"', data)
    if name == "ppt/slides/slide1.xml":
        return re.sub(rb'(r:embed|r:link)="rId2"', rb'\1="rId3"', data)
    return data


def _rewrite_zip(
    source: Path,
    destination: Path,
    transform: Callable[[str, bytes], bytes],
) -> Path:
    """Copy a package entry by entry, passing each through `transform`.

    Entry-by-entry rather than extract-and-rezip so that every part not under
    test is carried across byte-for-byte, and the entry order is preserved.
    """
    with zipfile.ZipFile(source) as src:
        entries = [(item, src.read(item.filename)) for item in src.infolist()]
    with zipfile.ZipFile(destination, "w", zipfile.ZIP_DEFLATED) as dst:
        for item, data in entries:
            dst.writestr(item.filename, transform(item.filename, data))
    return destination


# ---------------------------------------------------------------------------
# Registry
# ---------------------------------------------------------------------------

#: Fixture name -> builder. The conftest fixtures are generated from this, so
#: adding a deck here is all it takes to make it requestable in a test.
BUILDERS = {
    "empty": build_empty,
    "simple": build_simple,
    "picture": build_picture,
    "shared_picture": build_shared_picture,
    "notes": build_notes,
    "chart": build_chart,
    "two_charts": build_two_charts,
    "hyperlink": build_hyperlink,
    "gap_rids": build_gap_rids,
}


def build_all(out_dir: Path) -> dict[str, Path]:
    """Build every registered deck into `out_dir`."""
    out_dir.mkdir(parents=True, exist_ok=True)
    return {name: build(out_dir / f"{name}.pptx") for name, build in BUILDERS.items()}


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "--out",
        type=Path,
        default=None,
        help="destination directory (default: a fresh temp directory)",
    )
    args = parser.parse_args()

    out_dir = args.out or Path(tempfile.mkdtemp(prefix="pptx-plus-fixtures-"))
    for name, built in build_all(out_dir).items():
        print(f"{name:<16} {built}  ({built.stat().st_size:,} bytes)")


if __name__ == "__main__":
    main()

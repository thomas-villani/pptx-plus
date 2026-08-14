"""Reordering a slide -- SPEC §5.3, §5.5.

Every assertion about the result runs against a deck saved and reopened,
never the in-memory `Presentation` (SPEC §10.2).
"""

from __future__ import annotations

import pytest
from pptx import Presentation

from pptx_plus._testing import SavedPackage, assert_package_integrity, roundtrip, saved
from pptx_plus.core.oxml import xpath
from pptx_plus.core.sections import section_lst
from pptx_plus.slides import move_slide
from pptx_plus.slides.resolve import SlideIndexError, SlideNotFoundError


def _titles(prs: Presentation) -> list[str]:
    """The deck's running order, read back as slide titles."""
    return [slide.shapes.title.text for slide in prs.slides]


def _order(prs: Presentation) -> list[str]:
    """Running order after a save and reopen -- the artifact, not the model."""
    return _titles(roundtrip(prs))


@pytest.fixture
def prs(deck) -> Presentation:
    """Four titled slides: 'Slide 1' .. 'Slide 4'."""
    return Presentation(str(deck("simple4")))


# ---------------------------------------------------------------------------
# The index semantics table from SPEC §5.5
# ---------------------------------------------------------------------------


@pytest.mark.parametrize(
    ("index", "to_index", "expected"),
    [
        (0, 2, [2, 3, 1, 4]),
        (2, 2, [1, 2, 3, 4]),
        (3, 0, [4, 1, 2, 3]),
        (0, 3, [2, 3, 4, 1]),
        (1, 0, [2, 1, 3, 4]),
    ],
)
def test_the_documented_index_semantics(
    prs: Presentation, index: int, to_index: int, expected: list[int]
) -> None:
    """`to_index` is the position in the *resulting* deck. SPEC §5.5."""
    move_slide(prs, index, to_index)
    assert _order(prs) == [f"Slide {n}" for n in expected]


@pytest.mark.parametrize("index", [0, 1, 2, 3])
def test_moving_a_slide_to_its_own_index_is_a_no_op(prs: Presentation, index: int) -> None:
    """The property that pins the semantics down. SPEC §5.5."""
    before = _order(prs)
    move_slide(prs, index, index)
    assert _order(prs) == before


def test_a_negative_to_index_counts_from_the_end(prs: Presentation) -> None:
    move_slide(prs, 0, -1)
    assert _order(prs) == ["Slide 2", "Slide 3", "Slide 4", "Slide 1"]


def test_a_negative_source_index_is_accepted(prs: Presentation) -> None:
    move_slide(prs, -1, 0)
    assert _order(prs) == ["Slide 4", "Slide 1", "Slide 2", "Slide 3"]


def test_a_slide_object_may_be_given_instead_of_an_index(prs: Presentation) -> None:
    move_slide(prs, prs.slides[0], 2)
    assert _order(prs) == ["Slide 2", "Slide 3", "Slide 1", "Slide 4"]


def test_moving_is_reversible(prs: Presentation) -> None:
    before = _order(prs)
    move_slide(prs, 0, 3)
    move_slide(prs, 3, 0)
    assert _order(prs) == before


# ---------------------------------------------------------------------------
# Out of range raises rather than clamping
# ---------------------------------------------------------------------------


@pytest.mark.parametrize("to_index", [4, 99, -5])
def test_an_out_of_range_to_index_raises(prs: Presentation, to_index: int) -> None:
    """Diverges from `list.insert`, which clamps. SPEC §5.5."""
    with pytest.raises(SlideIndexError, match="to_index"):
        move_slide(prs, 0, to_index)


def test_the_error_names_the_valid_positions(prs: Presentation) -> None:
    with pytest.raises(SlideIndexError, match=r"valid positions are 0\.\.3"):
        move_slide(prs, 0, 9)


def test_an_out_of_range_to_index_changes_nothing(prs: Presentation) -> None:
    """A rejected call must not have half-applied."""
    before = _order(prs)
    with pytest.raises(SlideIndexError):
        move_slide(prs, 0, 9)
    assert _order(prs) == before


def test_an_out_of_range_source_index_raises(prs: Presentation) -> None:
    with pytest.raises(SlideIndexError, match="slide index"):
        move_slide(prs, 9, 0)


def test_a_foreign_slide_raises(prs: Presentation, deck) -> None:
    other = Presentation(str(deck("simple")))
    with pytest.raises(SlideNotFoundError):
        move_slide(prs, other.slides[0], 0)


def test_a_single_slide_deck_accepts_only_index_zero() -> None:
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[6])
    move_slide(presentation, 0, 0)
    with pytest.raises(SlideIndexError):
        move_slide(presentation, 0, 1)


# ---------------------------------------------------------------------------
# What a move must NOT change
# ---------------------------------------------------------------------------


def test_a_move_changes_no_parts(prs: Presentation) -> None:
    """Pure ordering: the package holds the same parts, in the same places."""
    before = set(SavedPackage(saved(prs).blob).partnames)
    move_slide(prs, 0, 3)
    assert set(saved(prs).partnames) == before


def test_a_move_changes_no_slide_ids(prs: Presentation) -> None:
    """A slide keeps its identity across a reorder; only its position changes."""
    before = sorted(int(v) for v in xpath(saved(prs).presentation, "./p:sldIdLst/p:sldId/@id"))
    move_slide(prs, 0, 3)
    after = sorted(int(v) for v in xpath(saved(prs).presentation, "./p:sldIdLst/p:sldId/@id"))
    assert after == before


def test_a_move_preserves_package_integrity(prs: Presentation) -> None:
    move_slide(prs, 0, 3)
    assert_package_integrity(saved(prs))


def test_a_moved_deck_survives_a_second_round_trip(prs: Presentation) -> None:
    """Reopening and re-saving the result must still be valid."""
    move_slide(prs, 3, 0)
    assert_package_integrity(saved(roundtrip(prs)))


# ---------------------------------------------------------------------------
# Sections follow the slide
# ---------------------------------------------------------------------------


def _sections(prs: Presentation) -> list[list[int]]:
    """Each section's slide ids, in order."""
    root = section_lst(prs)
    if root is None:
        return []
    return [
        [int(value) for value in xpath(id_list, "./p14:sldId/@id")]
        for id_list in xpath(root, "./p14:section/p14:sldIdLst")
    ]


@pytest.fixture
def sectioned(deck) -> Presentation:
    """Four slides in two sections of two."""
    return Presentation(str(deck("sections")))


def test_the_sectioned_fixture_starts_partitioned(sectioned: Presentation) -> None:
    assert [len(s) for s in _sections(sectioned)] == [2, 2]


@pytest.mark.parametrize("index", [0, 1, 2, 3])
def test_moving_a_slide_to_its_own_index_leaves_sections_alone(
    sectioned: Presentation, index: int
) -> None:
    """The boundary tie-break exists for exactly this. SPEC §4.7.

    Landing on a section boundary is ambiguous -- appending to one section and
    prepending to the next produce the same running order. Any tie-break other
    than "prefer the section it came from" silently re-sections the deck on a
    call that changed nothing.
    """
    before = _sections(sectioned)
    move_slide(sectioned, index, index)
    assert _sections(sectioned) == before


def test_a_slide_moved_across_a_boundary_changes_section(sectioned: Presentation) -> None:
    """Sections partition the running order, so crossing a boundary re-sections."""
    move_slide(sectioned, 0, 2)
    assert _sections(sectioned) == [[257], [258, 256, 259]]


def test_a_slide_moved_to_the_end_joins_the_last_section(sectioned: Presentation) -> None:
    move_slide(sectioned, 0, 3)
    assert _sections(sectioned) == [[257], [258, 259, 256]]


def test_a_move_within_one_section_does_not_cross(sectioned: Presentation) -> None:
    move_slide(sectioned, 0, 1)
    assert _sections(sectioned) == [[257, 256], [258, 259]]


def test_sections_stay_consistent_with_the_deck(sectioned: Presentation) -> None:
    """The flattened sections must equal the running order, always."""
    move_slide(sectioned, 3, 0)
    pkg = saved(sectioned)
    deck_ids = [int(v) for v in xpath(pkg.presentation, "./p:sldIdLst/p:sldId/@id")]
    assert [sld_id for section in _sections(sectioned) for sld_id in section] == deck_ids


def test_a_sectioned_move_preserves_package_integrity(sectioned: Presentation) -> None:
    move_slide(sectioned, 0, 3)
    assert_package_integrity(saved(sectioned))


def test_a_deck_without_sections_is_untouched(prs: Presentation) -> None:
    """No sections is the common case and must not grow an empty section list."""
    move_slide(prs, 0, 2)
    assert section_lst(roundtrip(prs)) is None


def test_a_custom_show_is_not_reordered_by_a_move(deck) -> None:
    """A custom show has its own running order, independent of the deck's."""
    presentation = Presentation(str(deck("custom_show")))
    before = [
        str(v) for v in xpath(saved(presentation).presentation, "./p:custShowLst//p:sld/@r:id")
    ]
    move_slide(presentation, 0, 2)
    after = [
        str(v) for v in xpath(saved(presentation).presentation, "./p:custShowLst//p:sld/@r:id")
    ]
    assert after == before

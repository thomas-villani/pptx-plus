"""Headless-render smoke tests -- SPEC §10.4.

**LibreOffice is a corruption detector, not a fidelity oracle.** It does not
run the SmartArt layout engine, its chart rendering differs from PowerPoint's,
and its font substitution is its own. So the only thing asserted here is that
a deck this library produced converts cleanly and yields the expected number of
pages. Claiming more would be a lie, and a comparison against a PowerPoint
reference render would fail for reasons that have nothing to do with this
library.

What it *does* catch is the class of package-level damage that a Python-side
assertion cannot see: a part the consumer cannot resolve, a content type it
rejects, a structure it refuses to open. Every deck reaching this point has
already passed the integrity battery, so a failure here means the battery is
missing an invariant -- which is the most valuable thing this tier can tell us.

Skipped unless `soffice` is on PATH, so a dev box without LibreOffice runs the
rest of the suite normally. CI runs it on the Linux leg.
"""

from __future__ import annotations

import re
import shutil
import subprocess
from pathlib import Path

import pytest
from pptx import Presentation

from pptx_plus.slides import delete_slide, duplicate_slide, move_slide

pytestmark = [
    pytest.mark.requires_libreoffice,
    pytest.mark.skipif(shutil.which("soffice") is None, reason="LibreOffice not on PATH"),
]

#: Matches a page object in a PDF's body. Crude, and sufficient: the decks here
#: are small and produced by one converter, so the count is exact. Adding a PDF
#: parser as a test dependency to count pages would not be a good trade.
_PDF_PAGE = re.compile(rb"/Type\s*/Page[^s]")


def _to_pdf(source: Path, out_dir: Path) -> Path:
    """Convert a deck to PDF with LibreOffice, returning the PDF's path."""
    result = subprocess.run(
        [
            "soffice",
            "--headless",
            "--norestore",
            "--convert-to",
            "pdf",
            "--outdir",
            str(out_dir),
            str(source),
        ],
        capture_output=True,
        timeout=180,
        check=False,
    )
    pdf = out_dir / f"{source.stem}.pdf"
    if not pdf.exists():
        raise AssertionError(
            f"LibreOffice produced no PDF for {source.name}.\n"
            f"exit={result.returncode}\n"
            f"stdout={result.stdout.decode(errors='replace')}\n"
            f"stderr={result.stderr.decode(errors='replace')}"
        )
    return pdf


def _page_count(pdf: Path) -> int:
    return len(_PDF_PAGE.findall(pdf.read_bytes()))


def _converted(prs: Presentation, tmp_path: Path, name: str = "deck") -> int:
    """Save, convert, and return the page count."""
    source = tmp_path / f"{name}.pptx"
    prs.save(source)
    return _page_count(_to_pdf(source, tmp_path))


# ---------------------------------------------------------------------------
# The fixtures convert as authored
# ---------------------------------------------------------------------------


@pytest.mark.parametrize(
    "name", ["simple", "simple4", "picture", "chart", "two_charts", "notes", "sections"]
)
def test_a_fixture_converts_cleanly(deck, tmp_path: Path, name: str) -> None:
    """The control: these decks were never touched by this library."""
    assert _page_count(_to_pdf(Path(deck(name)), tmp_path)) >= 1


# ---------------------------------------------------------------------------
# Decks this library produced
# ---------------------------------------------------------------------------


def test_a_duplicated_deck_converts_with_one_more_page(deck, tmp_path: Path) -> None:
    prs = Presentation(str(deck("simple4")))
    duplicate_slide(prs, 0)
    assert _converted(prs, tmp_path) == 5


def test_a_deleted_deck_converts_with_one_fewer_page(deck, tmp_path: Path) -> None:
    prs = Presentation(str(deck("simple4")))
    delete_slide(prs, 1)
    assert _converted(prs, tmp_path) == 3


def test_a_moved_deck_converts_with_the_same_page_count(deck, tmp_path: Path) -> None:
    prs = Presentation(str(deck("simple4")))
    move_slide(prs, 0, 3)
    assert _converted(prs, tmp_path) == 4


def test_a_duplicated_picture_slide_converts(deck, tmp_path: Path) -> None:
    """The image is shared, so the copy must still resolve it."""
    prs = Presentation(str(deck("picture")))
    duplicate_slide(prs, 0)
    assert _converted(prs, tmp_path) == 2


def test_a_duplicated_chart_slide_converts(deck, tmp_path: Path) -> None:
    """The copy owns its own chart part and embedded workbook."""
    prs = Presentation(str(deck("chart")))
    duplicate_slide(prs, 0)
    assert _converted(prs, tmp_path) == 2


def test_a_duplicated_two_chart_slide_converts(deck, tmp_path: Path) -> None:
    """The part-name reservation case, checked by an actual consumer.

    Without the reservation set both cloned charts land on one part name and
    one silently overwrites the other in the zip -- which a renderer notices
    and a namelist check might not, if the check looked at the wrong thing.
    """
    prs = Presentation(str(deck("two_charts")))
    duplicate_slide(prs, 0)
    assert _converted(prs, tmp_path) == 2


def test_a_gap_rids_duplicate_converts(deck, tmp_path: Path) -> None:
    """The relationship-remap regression, end to end through a renderer."""
    prs = Presentation(str(deck("gap_rids")))
    duplicate_slide(prs, 0)
    assert _converted(prs, tmp_path) == 2


def test_a_sectioned_deck_converts_after_every_verb(deck, tmp_path: Path) -> None:
    prs = Presentation(str(deck("sections")))
    duplicate_slide(prs, 0)
    move_slide(prs, 0, 3)
    delete_slide(prs, 1)
    assert _converted(prs, tmp_path) == 4


def test_an_emptied_deck_converts(deck, tmp_path: Path) -> None:
    """A zero-slide deck is schema-valid; check a consumer agrees."""
    prs = Presentation(str(deck("simple")))
    for _ in range(3):
        delete_slide(prs, 0)
    prs.save(tmp_path / "empty.pptx")
    _to_pdf(tmp_path / "empty.pptx", tmp_path)

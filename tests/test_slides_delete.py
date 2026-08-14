"""Deleting a slide -- SPEC §5.2, §5.6, §5.7.

Every assertion about what survives runs against a saved package, never the
in-memory `Presentation`. That is not a stylistic choice: after a delete the
object graph deliberately keeps the detached part, so an in-memory assertion
would fail on a perfectly correct delete and pass on the naive recipe that
leaves the slide in the file. SPEC §3.5.
"""

from __future__ import annotations

import contextlib

import pytest
from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

from pptx_plus._testing import (
    assert_in_package,
    assert_not_in_package,
    assert_package_integrity,
    roundtrip,
    saved,
)
from pptx_plus.core.oxml import part_root, xpath
from pptx_plus.core.sections import custom_show_lst, section_lst
from pptx_plus.slides import contains, delete_slide, move_slide
from pptx_plus.slides.resolve import SlideIndexError, SlideNotFoundError


def _titles(prs: Presentation) -> list[str]:
    return [slide.shapes.title.text for slide in roundtrip(prs).slides]


@pytest.fixture
def prs(deck) -> Presentation:
    return Presentation(str(deck("simple4")))


# ---------------------------------------------------------------------------
# The running order
# ---------------------------------------------------------------------------


def test_the_named_slide_is_the_one_removed(prs: Presentation) -> None:
    delete_slide(prs, 1)
    assert _titles(prs) == ["Slide 1", "Slide 3", "Slide 4"]


def test_the_deck_gets_shorter(prs: Presentation) -> None:
    delete_slide(prs, 1)
    assert len(roundtrip(prs).slides) == 3


def test_a_negative_index_counts_from_the_end(prs: Presentation) -> None:
    delete_slide(prs, -1)
    assert _titles(prs) == ["Slide 1", "Slide 2", "Slide 3"]


def test_a_slide_object_may_be_given_instead_of_an_index(prs: Presentation) -> None:
    delete_slide(prs, prs.slides[0])
    assert _titles(prs) == ["Slide 2", "Slide 3", "Slide 4"]


def test_deleting_the_first_slide(prs: Presentation) -> None:
    delete_slide(prs, 0)
    assert _titles(prs) == ["Slide 2", "Slide 3", "Slide 4"]


def test_deleting_the_last_slide_of_several(prs: Presentation) -> None:
    delete_slide(prs, 3)
    assert _titles(prs) == ["Slide 1", "Slide 2", "Slide 3"]


@pytest.mark.parametrize("index", [4, 99, -5])
def test_an_out_of_range_index_raises(prs: Presentation, index: int) -> None:
    with pytest.raises(SlideIndexError):
        delete_slide(prs, index)


def test_a_rejected_delete_changes_nothing(prs: Presentation) -> None:
    before = _titles(prs)
    with pytest.raises(SlideIndexError):
        delete_slide(prs, 9)
    assert _titles(prs) == before


# ---------------------------------------------------------------------------
# The empty deck is legal -- SPEC §5.2
# ---------------------------------------------------------------------------


def test_the_only_slide_of_a_deck_can_be_deleted(deck) -> None:
    """An empty `<p:sldIdLst/>` is schema-valid; PowerPoint opens a 0-slide deck."""
    presentation = Presentation(str(deck("picture")))
    delete_slide(presentation, 0)
    assert len(roundtrip(presentation).slides) == 0


def test_an_emptied_deck_is_still_a_valid_package(deck) -> None:
    presentation = Presentation(str(deck("simple")))
    for _ in range(3):
        delete_slide(presentation, 0)
    assert_package_integrity(saved(presentation))


def test_an_emptied_deck_can_be_reopened_and_used(deck) -> None:
    presentation = Presentation(str(deck("simple")))
    for _ in range(3):
        delete_slide(presentation, 0)
    reopened = roundtrip(presentation)
    reopened.slides.add_slide(reopened.slide_layouts[6])
    assert_package_integrity(saved(reopened))


# ---------------------------------------------------------------------------
# Idempotency contract -- SPEC §5.6
# ---------------------------------------------------------------------------


def test_deleting_the_same_slide_twice_raises(prs: Presentation) -> None:
    slide = prs.slides[1]
    delete_slide(prs, slide)
    with pytest.raises(SlideNotFoundError):
        delete_slide(prs, slide)


def test_suppress_key_error_gives_opt_in_idempotence(prs: Presentation) -> None:
    """The one line that buys idempotence without hiding it from everyone else."""
    slide = prs.slides[1]
    delete_slide(prs, slide)
    with contextlib.suppress(KeyError):
        delete_slide(prs, slide)
    assert _titles(prs) == ["Slide 1", "Slide 3", "Slide 4"]


def test_deleting_index_zero_twice_deletes_two_slides(prs: Presentation) -> None:
    """Index semantics, not a contradiction of the paragraph above. SPEC §5.6."""
    delete_slide(prs, 0)
    delete_slide(prs, 0)
    assert _titles(prs) == ["Slide 3", "Slide 4"]


# ---------------------------------------------------------------------------
# The stale Slide object -- SPEC §5.7
# ---------------------------------------------------------------------------


def test_a_deleted_slide_object_stays_readable(prs: Presentation) -> None:
    """Deletion detaches a part from the graph; it destroys nothing."""
    slide = prs.slides[1]
    delete_slide(prs, slide)
    assert slide.shapes.title.text == "Slide 2"


def test_a_deleted_slide_is_no_longer_contained(prs: Presentation) -> None:
    slide = prs.slides[1]
    delete_slide(prs, slide)
    assert contains(prs, slide) is False


def test_the_library_does_not_poison_the_deleted_object(prs: Presentation) -> None:
    """It is not our object. Mutating it to enforce a contract is a later bug."""
    slide = prs.slides[1]
    before = type(slide)
    delete_slide(prs, slide)
    assert type(slide) is before


# ---------------------------------------------------------------------------
# What leaves the package, and what does not
# ---------------------------------------------------------------------------


def test_the_slide_part_is_no_longer_written(prs: Presentation) -> None:
    """The whole difference from the naive recipe, in one assertion."""
    delete_slide(prs, 1)
    assert_not_in_package(saved(prs), "/ppt/slides/slide2.xml")


def test_the_surviving_slide_parts_are_still_written(prs: Presentation) -> None:
    delete_slide(prs, 1)
    pkg = saved(prs)
    for partname in ("/ppt/slides/slide1.xml", "/ppt/slides/slide3.xml"):
        assert_in_package(pkg, partname)


def test_an_unshared_image_is_collected(deck) -> None:
    """Dropping the relationship is the whole of collecting it. SPEC §3.5."""
    presentation = Presentation(str(deck("picture")))
    assert_in_package(saved(presentation), "/ppt/media/image1.gif")
    delete_slide(presentation, 0)
    assert_not_in_package(saved(presentation), "/ppt/media/image1.gif")


def test_a_shared_image_survives(deck) -> None:
    """Still reachable from the other slide, so still written -- for free."""
    presentation = Presentation(str(deck("shared_picture")))
    delete_slide(presentation, 0)
    assert_in_package(saved(presentation), "/ppt/media/image1.gif")


def test_the_other_slide_can_still_resolve_the_shared_image(deck) -> None:
    presentation = Presentation(str(deck("shared_picture")))
    delete_slide(presentation, 0)
    assert_package_integrity(saved(presentation))


def test_a_notes_slide_goes_with_its_slide(deck) -> None:
    presentation = Presentation(str(deck("notes")))
    assert_in_package(saved(presentation), "/ppt/notesSlides/notesSlide1.xml")
    delete_slide(presentation, 1)
    assert_not_in_package(saved(presentation), "/ppt/notesSlides/notesSlide1.xml")


def test_the_notes_master_survives_its_last_notes_slide(deck) -> None:
    """The master hangs off the presentation part, not off any slide."""
    presentation = Presentation(str(deck("notes")))
    delete_slide(presentation, 1)
    pkg = saved(presentation)
    masters = pkg.targets(pkg.presentation_partname, [RT.NOTES_MASTER])
    assert len(masters) == 1


def test_layouts_and_masters_survive_deleting_every_slide(deck) -> None:
    """They are reached from the master's own list, never from a slide."""
    presentation = Presentation(str(deck("simple")))
    for _ in range(3):
        delete_slide(presentation, 0)
    pkg = saved(presentation)
    assert_in_package(pkg, "/ppt/slideMasters/slideMaster1.xml")
    assert_in_package(pkg, "/ppt/slideLayouts/slideLayout1.xml")


def test_the_package_actually_shrinks(deck) -> None:
    """Not merely "the part is absent" -- the file a user gets is smaller."""
    presentation = Presentation(str(deck("picture")))
    before = len(saved(presentation).blob)
    delete_slide(presentation, 0)
    assert len(saved(presentation).blob) < before


def test_slide_part_names_are_not_renumbered(prs: Presentation) -> None:
    """A gap is normal and harmless. SPEC §5.8."""
    delete_slide(prs, 1)
    assert_in_package(saved(prs), "/ppt/slides/slide3.xml")


# ---------------------------------------------------------------------------
# Sections and custom shows -- the corruption vectors
# ---------------------------------------------------------------------------


def _sections(prs: Presentation) -> list[list[int]]:
    root = section_lst(prs)
    if root is None:
        return []
    return [
        [int(value) for value in xpath(id_list, "./p14:sldId/@id")]
        for id_list in xpath(root, "./p14:section/p14:sldIdLst")
    ]


def test_a_deleted_slide_leaves_its_section(deck) -> None:
    presentation = Presentation(str(deck("sections")))
    delete_slide(presentation, 0)
    assert _sections(presentation) == [[257], [258, 259]]


def test_a_sectioned_delete_is_section_consistent(deck) -> None:
    """The invariant a naive delete violates, and the reason §4.7 exists."""
    presentation = Presentation(str(deck("sections")))
    delete_slide(presentation, 1)
    assert_package_integrity(saved(presentation))


def test_an_emptied_section_is_left_in_place(deck) -> None:
    """A named thing the user created; an empty one is schema-valid."""
    presentation = Presentation(str(deck("sections")))
    delete_slide(presentation, 0)
    delete_slide(presentation, 0)
    assert _sections(presentation) == [[], [258, 259]]


def test_a_deleted_slide_leaves_every_custom_show(deck) -> None:
    presentation = Presentation(str(deck("custom_show")))
    delete_slide(presentation, 0)
    shows = custom_show_lst(presentation)
    assert shows is not None
    assert len(xpath(shows, "./p:custShow/p:sldLst/p:sld")) == 1


def test_a_custom_show_delete_is_consistent(deck) -> None:
    """The case where `XmlPart.drop_rel` would have silently no-opped."""
    presentation = Presentation(str(deck("custom_show")))
    delete_slide(presentation, 0)
    assert_package_integrity(saved(presentation))


def test_a_slide_in_a_custom_show_really_is_referenced_twice(deck) -> None:
    """Pins the premise: this is what makes `drop_rel` refuse. SPEC §4.6."""
    presentation = Presentation(str(deck("custom_show")))
    root = part_root(presentation.part)
    r_id = str(xpath(root, "./p:sldIdLst/p:sldId[1]/@r:id")[0])
    assert len(xpath(root, "//@r:id[.=$rid]", rid=r_id)) == 2


def test_deleting_a_slide_not_in_the_custom_show_leaves_it_alone(deck) -> None:
    presentation = Presentation(str(deck("custom_show")))
    delete_slide(presentation, 1)
    shows = custom_show_lst(presentation)
    assert shows is not None
    assert len(xpath(shows, "./p:custShow/p:sldLst/p:sld")) == 2


# ---------------------------------------------------------------------------
# Integrity, and composition with move
# ---------------------------------------------------------------------------


@pytest.mark.parametrize("index", [0, 1, 2, 3])
def test_every_delete_leaves_a_valid_package(prs: Presentation, index: int) -> None:
    delete_slide(prs, index)
    assert_package_integrity(saved(prs))


def test_a_deleted_deck_survives_a_second_round_trip(prs: Presentation) -> None:
    delete_slide(prs, 1)
    assert_package_integrity(saved(roundtrip(prs)))


def test_delete_and_move_compose(prs: Presentation) -> None:
    delete_slide(prs, 1)
    move_slide(prs, 0, 2)
    assert _titles(prs) == ["Slide 3", "Slide 4", "Slide 1"]


def test_deleting_a_shared_image_deck_repeatedly_stays_valid(deck) -> None:
    presentation = Presentation(str(deck("shared_picture")))
    delete_slide(presentation, 0)
    delete_slide(presentation, 0)
    assert_package_integrity(saved(presentation))


# ---------------------------------------------------------------------------
# Known limitation -- slide-jump hyperlinks
# ---------------------------------------------------------------------------


def test_a_slide_jump_target_keeps_the_deleted_part_reachable(deck) -> None:
    """Documented v0.1 limitation, asserted so it cannot change unnoticed.

    Slide 2 holds an `RT.SLIDE` relationship to slide 1, so deleting slide 1
    leaves its part reachable and therefore still written -- alive in the file
    with no `p:sldId`. PowerPoint tolerates it, but the link goes nowhere;
    real PowerPoint would remove the hyperlink. `scrub_links=True` is the v0.2
    fix (ROADMAP).

    The package stays *valid* — the battery passes — which is exactly why this
    needs an explicit test rather than trusting the battery to notice.
    """
    presentation = Presentation(str(deck("slide_jump")))
    delete_slide(presentation, 0)
    pkg = saved(presentation)
    assert_package_integrity(pkg)
    assert_in_package(pkg, "/ppt/slides/slide1.xml")


def test_a_slide_jump_deck_has_no_slide_id_for_the_orphan(deck) -> None:
    presentation = Presentation(str(deck("slide_jump")))
    delete_slide(presentation, 0)
    assert len(xpath(saved(presentation).presentation, "./p:sldIdLst/p:sldId")) == 1

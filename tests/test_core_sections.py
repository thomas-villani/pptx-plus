"""Sections and custom shows -- SPEC §4.7.

The structures python-pptx does not model. They survive its round trips only
because unrecognized XML is preserved verbatim, which is what makes a stale
entry so easy to leave behind and so invisible when you do.
"""

from __future__ import annotations

import pytest
from pptx import Presentation

from pptx_plus._testing import assert_sections_consistent, roundtrip, saved
from pptx_plus.core.oxml import sld_id_lst, xpath
from pptx_plus.core.sections import (
    custom_show_lst,
    reorder_slide,
    scrub_slide,
    section_lst,
)


def _section_ids(prs: Presentation) -> list[list[int]]:
    root = section_lst(prs)
    if root is None:
        return []
    return [
        [int(value) for value in xpath(id_list, "./p14:sldId/@id")]
        for id_list in xpath(root, "./p14:section/p14:sldIdLst")
    ]


def _show_rids(prs: Presentation) -> list[str]:
    shows = custom_show_lst(prs)
    if shows is None:
        return []
    return [str(value) for value in xpath(shows, "./p:custShow/p:sldLst/p:sld/@r:id")]


@pytest.fixture
def sectioned(deck) -> Presentation:
    return Presentation(str(deck("sections")))


@pytest.fixture
def with_show(deck) -> Presentation:
    return Presentation(str(deck("custom_show")))


# ---------------------------------------------------------------------------
# Locating the structures
# ---------------------------------------------------------------------------


def test_section_lst_is_found_when_present(sectioned: Presentation) -> None:
    assert section_lst(sectioned) is not None


def test_section_lst_is_none_when_absent(deck) -> None:
    assert section_lst(Presentation(str(deck("simple")))) is None


def test_section_lst_does_not_create_the_element(deck) -> None:
    """A deck without sections must not grow an empty section list. SPEC §9.8."""
    prs = Presentation(str(deck("simple")))
    section_lst(prs)
    assert b"sectionLst" not in saved(prs).blob


def test_custom_show_lst_is_found_when_present(with_show: Presentation) -> None:
    assert custom_show_lst(with_show) is not None


def test_custom_show_lst_is_none_when_absent(deck) -> None:
    assert custom_show_lst(Presentation(str(deck("simple")))) is None


def test_sections_survive_a_round_trip(sectioned: Presentation) -> None:
    """The premise the whole module rests on: python-pptx preserves them."""
    assert _section_ids(roundtrip(sectioned)) == [[256, 257], [258, 259]]


def test_custom_shows_survive_a_round_trip(with_show: Presentation) -> None:
    assert len(_show_rids(roundtrip(with_show))) == 2


# ---------------------------------------------------------------------------
# scrub_slide
# ---------------------------------------------------------------------------


def test_scrub_removes_the_section_entry(sectioned: Presentation) -> None:
    scrub_slide(sectioned, slide_id=256, r_id=sld_id_lst(sectioned)[0].rId)
    assert _section_ids(sectioned) == [[257], [258, 259]]


def test_scrub_removes_the_custom_show_entry(with_show: Presentation) -> None:
    entry = sld_id_lst(with_show)[0]
    scrub_slide(with_show, slide_id=int(entry.get("id")), r_id=entry.rId)
    assert len(_show_rids(with_show)) == 1


def test_scrub_leaves_other_entries_alone(with_show: Presentation) -> None:
    entries = list(sld_id_lst(with_show))
    survivor = entries[2].rId
    scrub_slide(with_show, slide_id=int(entries[0].get("id")), r_id=entries[0].rId)
    assert _show_rids(with_show) == [survivor]


def test_scrub_leaves_an_emptied_section_in_place(deck) -> None:
    """An empty section is a named thing the user made, and schema-valid.

    Removing it would be a second, unrequested edit -- and PowerPoint keeps it.
    """
    prs = Presentation(str(deck("sections")))
    for entry in list(sld_id_lst(prs))[:2]:
        scrub_slide(prs, slide_id=int(entry.get("id")), r_id=entry.rId)
    assert _section_ids(prs) == [[], [258, 259]]


def test_scrub_is_a_no_op_on_a_deck_with_neither(deck) -> None:
    prs = Presentation(str(deck("simple")))
    entry = sld_id_lst(prs)[0]
    scrub_slide(prs, slide_id=int(entry.get("id")), r_id=entry.rId)
    assert _section_ids(prs) == []


def test_scrub_of_an_absent_slide_changes_nothing(sectioned: Presentation) -> None:
    before = _section_ids(sectioned)
    scrub_slide(sectioned, slide_id=9999, r_id="rId9999")
    assert _section_ids(sectioned) == before


def test_a_scrubbed_deck_is_section_consistent(sectioned: Presentation) -> None:
    """The end-to-end point: scrubbing is what keeps I6 satisfied."""
    lst = sld_id_lst(sectioned)
    entry = lst[1]
    scrub_slide(sectioned, slide_id=int(entry.get("id")), r_id=entry.rId)
    sectioned.part.rels.pop(entry.rId)
    lst.remove(entry)
    assert_sections_consistent(saved(sectioned))


# ---------------------------------------------------------------------------
# reorder_slide
# ---------------------------------------------------------------------------


@pytest.mark.parametrize(("slide_id", "index"), [(256, 0), (257, 1), (258, 2), (259, 3)])
def test_reordering_to_the_same_position_is_a_no_op(
    sectioned: Presentation, slide_id: int, index: int
) -> None:
    """The boundary tie-break: a slide that did not move keeps its section."""
    before = _section_ids(sectioned)
    reorder_slide(sectioned, slide_id=slide_id, to_index=index)
    assert _section_ids(sectioned) == before


def test_reordering_across_a_boundary_moves_the_entry(sectioned: Presentation) -> None:
    """The moved slide changes section. Section 1 shrinks, section 2 grows."""
    reorder_slide(sectioned, slide_id=256, to_index=2)
    assert _section_ids(sectioned) == [[257], [258, 256, 259]]


def test_reordering_changes_no_other_slides_section(sectioned: Presentation) -> None:
    """The property that decides the ambiguous cases: exactly one membership moves.

    The rejected alternative holds each section's *size* fixed and lets the
    boundary stay at a fixed position, which re-sections a bystander -- moving
    slide 1 into section 2 would drag slide 3 back into section 1. Nobody
    dragging a slide in the slide sorter expects a different slide to change
    section as a result.
    """
    before = {
        sld_id: number
        for number, section in enumerate(_section_ids(sectioned))
        for sld_id in section
    }
    reorder_slide(sectioned, slide_id=256, to_index=2)
    after = {
        sld_id: number
        for number, section in enumerate(_section_ids(sectioned))
        for sld_id in section
    }
    assert {k: v for k, v in after.items() if before[k] != v} == {256: 1}


def test_reordering_to_the_last_position_joins_the_last_section(
    sectioned: Presentation,
) -> None:
    reorder_slide(sectioned, slide_id=256, to_index=3)
    assert _section_ids(sectioned) == [[257], [258, 259, 256]]


def test_reordering_to_the_first_position_joins_the_first_section(
    sectioned: Presentation,
) -> None:
    reorder_slide(sectioned, slide_id=259, to_index=0)
    assert _section_ids(sectioned) == [[259, 256, 257], [258]]


def test_reordering_preserves_the_total_entry_count(sectioned: Presentation) -> None:
    reorder_slide(sectioned, slide_id=256, to_index=3)
    assert sum(len(section) for section in _section_ids(sectioned)) == 4


def test_reordering_is_a_no_op_without_sections(deck) -> None:
    prs = Presentation(str(deck("simple")))
    reorder_slide(prs, slide_id=256, to_index=2)
    assert section_lst(prs) is None


def test_reordering_an_unsectioned_slide_changes_nothing(sectioned: Presentation) -> None:
    """A slide missing from the sections is an anomaly to leave alone.

    Inventing an entry would be a change the caller did not ask for, and it is
    not this function's business to repair a deck it was asked to reorder.
    """
    before = _section_ids(sectioned)
    reorder_slide(sectioned, slide_id=9999, to_index=0)
    assert _section_ids(sectioned) == before


def test_reordering_when_sections_do_not_cover_the_deck(sectioned: Presentation) -> None:
    """A deck can arrive with sections that omit slides; do not invent placement.

    The entry goes back where it was rather than landing somewhere derived
    from an index the sections cannot interpret.
    """
    root = section_lst(sectioned)
    assert root is not None
    for entry in xpath(root, "./p14:section[2]/p14:sldIdLst/p14:sldId"):
        entry.getparent().remove(entry)
    reorder_slide(sectioned, slide_id=256, to_index=3)
    assert _section_ids(sectioned) == [[257, 256], []]

"""The committed PowerPoint-authored samples -- SPEC §10.6.

Everything else in the suite runs on decks python-pptx built. That is a real
limit: python-pptx cannot author SmartArt at all, and generated decks have no
sections, no transitions, and no animations. A library can pass a suite built
entirely from its own output and still be wrong about the format.

These decks were authored by PowerPoint. Provenance is in
`tests/fixtures/pptx_samples/README.md`.

They earned their place immediately: `smartart.pptx` exposed two defects that
every generated fixture had hidden, both recorded in `test_the_diagram_data_
part_has_no_relationships_of_its_own` and the tests below it.
"""

from __future__ import annotations

import re
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

from pptx_plus._testing import (
    SavedPackage,
    assert_in_package,
    assert_not_in_package,
    assert_package_integrity,
    roundtrip,
    saved,
)
from pptx_plus.core.reltypes import DIAGRAM_RELTYPES, RT_DIAGRAM_DRAWING
from pptx_plus.core.sections import section_lst
from pptx_plus.slides import delete_slide, duplicate_slide, move_slide

SAMPLES = ["smartart", "real_world"]

_REL_ID = re.compile(rb'relId="([^"]*)"')


@pytest.fixture
def sample(sample_dir: Path):
    def load(name: str) -> Presentation:
        path = sample_dir / f"{name}.pptx"
        assert path.exists(), f"missing committed sample {path}; see its README"
        return Presentation(str(path))

    return load


def _pkg(sample_dir: Path, name: str) -> SavedPackage:
    return SavedPackage((sample_dir / f"{name}.pptx").read_bytes())


def _diagram_parts(pkg: SavedPackage, slide: str) -> dict[str, str]:
    """Map diagram reltype -> partname for one slide."""
    return {
        rel.reltype: rel.partname
        for rel in pkg.rels(slide).values()
        if rel.reltype in DIAGRAM_RELTYPES and rel.partname
    }


# ---------------------------------------------------------------------------
# The harness passes on real PowerPoint output -- SPEC §10.5
# ---------------------------------------------------------------------------


@pytest.mark.parametrize("name", SAMPLES)
def test_a_sample_passes_the_battery_unmodified(sample_dir: Path, name: str) -> None:
    """Step 1 of grading the harness, against input it did not produce.

    This is the assertion that caught the harness reporting every
    PowerPoint-authored diagram as damaged: `dsp:dataModelExt/@relId` resolves
    against the *referring* part, and the harness was checking it against the
    containing part's own (empty) relationship set.
    """
    assert_package_integrity(_pkg(sample_dir, name))


@pytest.mark.parametrize("name", SAMPLES)
def test_a_sample_survives_an_untouched_round_trip(sample: object, name: str) -> None:
    """python-pptx re-serializes what it parsed; the result must still be valid."""
    assert_package_integrity(saved(sample(name)))  # type: ignore[operator]


# ---------------------------------------------------------------------------
# What a real SmartArt deck actually looks like
# ---------------------------------------------------------------------------


def test_a_smartart_slide_has_all_five_diagram_parts(sample_dir: Path) -> None:
    """Four from ECMA-376 plus the Microsoft drawing-cache extension.

    python-pptx knows four of these reltypes and not `diagramDrawing`, which
    is why `RT_DIAGRAM_DRAWING` is defined by hand.
    """
    pkg = _pkg(sample_dir, "smartart")
    assert set(_diagram_parts(pkg, pkg.slide_partnames[0])) == set(DIAGRAM_RELTYPES)


def test_the_diagram_data_part_has_no_relationships_of_its_own(sample_dir: Path) -> None:
    """The premise behind the parent-scope rule, pinned against a real deck.

    Every other relationship reference in OOXML is scoped to the part that
    contains it (SPEC §3.2). This one is not, and the reason is visible here:
    the data part has no `.rels` at all, so its `relId` *cannot* be
    self-scoped.
    """
    pkg = _pkg(sample_dir, "smartart")
    data = _diagram_parts(pkg, pkg.slide_partnames[0])[RT.DIAGRAM_DATA]
    assert pkg.rels(data) == {}


def test_the_diagram_data_part_carries_an_unqualified_rel_id(sample_dir: Path) -> None:
    """The one relationship id in a real deck outside the `r:` namespace."""
    pkg = _pkg(sample_dir, "smartart")
    data = _diagram_parts(pkg, pkg.slide_partnames[0])[RT.DIAGRAM_DATA]
    assert _REL_ID.findall(pkg.blob_of(data))


def test_the_unqualified_rel_id_names_the_drawing_part(sample_dir: Path) -> None:
    """And it resolves through the *slide's* relationships, not its own."""
    pkg = _pkg(sample_dir, "smartart")
    slide = pkg.slide_partnames[0]
    parts = _diagram_parts(pkg, slide)
    rel_id = _REL_ID.findall(pkg.blob_of(parts[RT.DIAGRAM_DATA]))[0].decode()
    assert pkg.rels(slide)[rel_id].partname == parts[RT_DIAGRAM_DRAWING]


# ---------------------------------------------------------------------------
# Duplicating SmartArt
# ---------------------------------------------------------------------------


def test_duplicating_a_smartart_slide_is_valid(sample: object) -> None:
    prs = sample("smartart")  # type: ignore[operator]
    duplicate_slide(prs, 0)
    assert_package_integrity(saved(prs))


def test_the_copy_gets_its_own_five_diagram_parts(sample: object) -> None:
    """Every diagram part is a mutable definition, so none may be shared."""
    prs = sample("smartart")  # type: ignore[operator]
    duplicate_slide(prs, 0)
    pkg = saved(prs)
    first, second = pkg.slide_partnames
    original = _diagram_parts(pkg, first)
    copy = _diagram_parts(pkg, second)
    assert set(original.values()).isdisjoint(copy.values())


def test_the_copys_diagram_points_at_its_own_drawing(sample: object) -> None:
    """**The defect a generated fixture could never have caught.**

    python-pptx has no model class for a diagram data part, so it loads as an
    opaque blob with no element tree -- and the ordinary rewrite pass skipped
    it entirely. The copy kept the source's `relId`, so its diagram pointed at
    the *original's* drawing cache. No exception, no repair prompt: PowerPoint
    recomputes the drawing on open and looks fine, while every other renderer
    draws the wrong thing or nothing.
    """
    prs = sample("smartart")  # type: ignore[operator]
    duplicate_slide(prs, 0)
    pkg = saved(prs)
    for slide in pkg.slide_partnames:
        parts = _diagram_parts(pkg, slide)
        rel_id = _REL_ID.findall(pkg.blob_of(parts[RT.DIAGRAM_DATA]))[0].decode()
        assert pkg.rels(slide)[rel_id].partname == parts[RT_DIAGRAM_DRAWING], slide


@pytest.mark.parametrize(
    "reltype", [RT.DIAGRAM_LAYOUT, RT.DIAGRAM_COLORS, RT.DIAGRAM_QUICK_STYLE, RT_DIAGRAM_DRAWING]
)
def test_the_other_diagram_parts_are_copied_byte_for_byte(sample: object, reltype: str) -> None:
    """SPEC §8.1. Only the data part is reparsed, and only because it must be."""
    prs = sample("smartart")  # type: ignore[operator]
    duplicate_slide(prs, 0)
    pkg = saved(prs)
    first, second = pkg.slide_partnames
    assert pkg.blob_of(_diagram_parts(pkg, first)[reltype]) == pkg.blob_of(
        _diagram_parts(pkg, second)[reltype]
    )


def test_duplicating_smartart_twice_stays_valid(sample: object) -> None:
    prs = sample("smartart")  # type: ignore[operator]
    duplicate_slide(prs, 0)
    duplicate_slide(prs, 0)
    assert_package_integrity(saved(prs))


def test_each_smartart_copy_owns_a_distinct_data_part(sample: object) -> None:
    prs = sample("smartart")  # type: ignore[operator]
    duplicate_slide(prs, 0)
    duplicate_slide(prs, 0)
    pkg = saved(prs)
    data = [_diagram_parts(pkg, s)[RT.DIAGRAM_DATA] for s in pkg.slide_partnames]
    assert len(set(data)) == 3


# ---------------------------------------------------------------------------
# Deleting SmartArt
# ---------------------------------------------------------------------------


def test_deleting_a_smartart_slide_collects_every_diagram_part(sample: object) -> None:
    prs = sample("smartart")  # type: ignore[operator]
    before = _diagram_parts(saved(prs), saved(prs).slide_partnames[0])
    delete_slide(prs, 0)
    pkg = saved(prs)
    for partname in before.values():
        assert_not_in_package(pkg, partname)


def test_deleting_one_smartart_copy_leaves_the_others(sample: object) -> None:
    """The copies own separate parts, so collecting one must not take another's."""
    prs = sample("smartart")  # type: ignore[operator]
    copy = duplicate_slide(prs, 0)
    survivor = _diagram_parts(saved(prs), saved(prs).slide_partnames[0])
    delete_slide(prs, copy)
    pkg = saved(prs)
    assert_package_integrity(pkg)
    for partname in survivor.values():
        assert_in_package(pkg, partname)


# ---------------------------------------------------------------------------
# The real-world deck: PowerPoint sections, notes, a slide-jump link
# ---------------------------------------------------------------------------


def test_the_real_world_deck_has_powerpoint_authored_sections(sample: object) -> None:
    """Not the hand-built approximation the generated fixtures use."""
    prs = sample("real_world")  # type: ignore[operator]
    assert section_lst(prs) is not None


def test_powerpoint_sections_survive_a_round_trip(sample: object) -> None:
    prs = sample("real_world")  # type: ignore[operator]
    assert section_lst(roundtrip(prs)) is not None


@pytest.mark.parametrize("verb", ["move", "duplicate", "delete"])
def test_every_verb_keeps_the_real_deck_valid(sample: object, verb: str) -> None:
    prs = sample("real_world")  # type: ignore[operator]
    if verb == "move":
        move_slide(prs, 0, 2)
    elif verb == "duplicate":
        duplicate_slide(prs, 0)
    else:
        delete_slide(prs, 1)
    assert_package_integrity(saved(prs))


def test_the_real_deck_survives_all_three_verbs_in_sequence(sample: object) -> None:
    prs = sample("real_world")  # type: ignore[operator]
    duplicate_slide(prs, 0)
    move_slide(prs, 0, 3)
    delete_slide(prs, 1)
    assert_package_integrity(saved(roundtrip(prs)))


def test_speaker_notes_authored_in_powerpoint_are_copied(sample: object) -> None:
    prs = sample("real_world")  # type: ignore[operator]
    source = next(s for s in prs.slides if s.has_notes_slide)
    copy = duplicate_slide(prs, source)
    assert copy.notes_slide.notes_text_frame.text == (source.notes_slide.notes_text_frame.text)


def test_a_powerpoint_slide_jump_is_not_duplicated(sample: object) -> None:
    """Authored by PowerPoint rather than by python-pptx's click_action."""
    prs = sample("real_world")  # type: ignore[operator]
    before = len(saved(prs).slide_partnames)
    duplicate_slide(prs, 2)
    assert len(saved(prs).slide_partnames) == before + 1

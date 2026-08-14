"""Part-name allocation, cloning, and relationship removal — SPEC §4.6."""

from __future__ import annotations

import io

import pytest
from pptx import Presentation
from pptx.opc.constants import CONTENT_TYPE as CT
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.util import Inches

from pptx_plus.core.parts import (
    UnclonablePartError,
    allocate_partname,
    clone_part,
    drop_relationship,
    partname_template_for,
)

# A 1x1 GIF — the smallest thing python-pptx will accept as a picture.
GIF = (
    b"GIF89a\x01\x00\x01\x00\x80\x00\x00\x00\x00\x00\xff\xff\xff!"
    b"\xf9\x04\x01\x00\x00\x00\x00,\x00\x00\x00\x00\x01\x00\x01\x00\x00\x02\x02D\x01\x00;"
)


def _deck_with_picture() -> Presentation:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(io.BytesIO(GIF), Inches(1), Inches(1))
    return prs


# --------------------------------------------------------------------------
# partname_template_for
# --------------------------------------------------------------------------


def test_template_comes_from_the_table_for_a_known_type() -> None:
    prs = Presentation()
    part = prs.slides.add_slide(prs.slide_layouts[6]).part
    assert partname_template_for(part) == "/ppt/slides/slide%d.xml"


def test_template_is_derived_for_an_unknown_type() -> None:
    """An unrecognized part still clones, landing beside its original.

    "Unknown to us" describes most of what makes a real deck interesting, so
    refusing to clone unknown parts would defeat the preservation guarantee.
    """
    prs = _deck_with_picture()
    image = next(r.target_part for r in prs.slides[0].part.rels.values() if r.reltype == RT.IMAGE)
    assert partname_template_for(image) == "/ppt/media/image%d.gif"


def test_unnumbered_unknown_part_is_refused(monkeypatch: pytest.MonkeyPatch) -> None:
    prs = Presentation()
    part = prs.slides.add_slide(prs.slide_layouts[6]).part
    monkeypatch.setattr(type(part), "content_type", "application/x-unknown", raising=False)
    monkeypatch.setattr(type(part), "partname", "/ppt/odd/thing.bin", raising=False)
    with pytest.raises(UnclonablePartError, match="no numeric component"):
        partname_template_for(part)


# --------------------------------------------------------------------------
# allocate_partname
# --------------------------------------------------------------------------


def test_allocate_skips_a_name_already_in_the_package() -> None:
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    allocated = allocate_partname(prs.part.package, "/ppt/slides/slide%d.xml", set())
    assert str(allocated) == "/ppt/slides/slide2.xml"


def test_allocate_skips_a_reserved_name() -> None:
    """The trap: `next_partname` cannot see a constructed-but-unattached part.

    Without the reservation set, cloning a slide with two charts allocates the
    same name twice and the second zip entry silently overwrites the first.
    """
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    reserved = {"/ppt/slides/slide2.xml"}
    allocated = allocate_partname(prs.part.package, "/ppt/slides/slide%d.xml", reserved)
    assert str(allocated) == "/ppt/slides/slide3.xml"


def test_allocate_records_its_own_reservation() -> None:
    prs = Presentation()
    reserved: set[str] = set()
    first = allocate_partname(prs.part.package, "/ppt/charts/chart%d.xml", reserved)
    second = allocate_partname(prs.part.package, "/ppt/charts/chart%d.xml", reserved)
    assert first != second


def test_repeated_allocation_never_collides() -> None:
    prs = Presentation()
    reserved: set[str] = set()
    names = [
        str(allocate_partname(prs.part.package, "/ppt/slides/slide%d.xml", reserved))
        for _ in range(5)
    ]
    assert len(set(names)) == 5


# --------------------------------------------------------------------------
# clone_part
# --------------------------------------------------------------------------


def test_clone_gets_a_fresh_partname() -> None:
    prs = Presentation()
    src = prs.slides.add_slide(prs.slide_layouts[6]).part
    clone = clone_part(src, into=src.package, reserved=set())
    assert clone.partname != src.partname


def test_clone_keeps_the_content_type() -> None:
    prs = Presentation()
    src = prs.slides.add_slide(prs.slide_layouts[6]).part
    clone = clone_part(src, into=src.package, reserved=set())
    assert clone.content_type == CT.PML_SLIDE


def test_clone_keeps_the_class() -> None:
    """`type(src).load(...)` — uniform across every part class."""
    prs = Presentation()
    src = prs.slides.add_slide(prs.slide_layouts[6]).part
    clone = clone_part(src, into=src.package, reserved=set())
    assert type(clone) is type(src)


def test_clone_has_no_relationships() -> None:
    """Minting them is the caller's job — the ids are what the XML rewrite needs."""
    prs = Presentation()
    src = prs.slides.add_slide(prs.slide_layouts[6]).part
    assert len(clone_part(src, into=src.package, reserved=set()).rels) == 0


def test_blob_backed_clone_is_byte_identical() -> None:
    """SPEC §8.1's preservation guarantee, at its narrowest.

    A blob-backed part is copied as bytes with no reparse, which is why
    SmartArt definition parts and embedded media survive a duplicate exactly
    rather than approximately. An image part stands in here for all of them —
    the diagram parts that make this matter cannot be authored by python-pptx.
    """
    prs = _deck_with_picture()
    image = next(r.target_part for r in prs.slides[0].part.rels.values() if r.reltype == RT.IMAGE)
    clone = clone_part(image, into=image.package, reserved=set())
    assert clone.blob == image.blob


def test_clone_is_independent_of_its_source() -> None:
    prs = Presentation()
    src = prs.slides.add_slide(prs.slide_layouts[6]).part
    clone = clone_part(src, into=src.package, reserved=set())
    assert clone is not src


# --------------------------------------------------------------------------
# drop_relationship
# --------------------------------------------------------------------------


def test_drop_relationship_removes_it() -> None:
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    rid = prs.slides._sldIdLst[0].rId
    drop_relationship(prs.part, rid)
    assert rid not in prs.part.rels


def test_drop_relationship_raises_on_an_unknown_id() -> None:
    with pytest.raises(KeyError):
        drop_relationship(Presentation().part, "rId9999")


def test_drop_relationship_is_unconditional() -> None:
    """The reason this exists instead of `XmlPart.drop_rel`.

    `drop_rel` refuses when the id is referenced twice or more in the part's
    XML. Here a second reference is added by hand — the shape a custom show
    produces — and `drop_rel` is shown to no-op on it while this function
    does not.
    """
    from pptx_plus.core.oxml import el, part_root, sld_id_lst

    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    rid = sld_id_lst(prs)[0].rId

    # A second `@r:id` reference to the same slide, as `p:custShowLst` gives.
    part_root(prs.part).append(el("p:custShowLst", **{"r:id": rid}))

    prs.part.drop_rel(rid)
    assert rid in prs.part.rels, "drop_rel unexpectedly removed a twice-referenced rel"

    drop_relationship(prs.part, rid)
    assert rid not in prs.part.rels

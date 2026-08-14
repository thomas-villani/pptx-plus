"""Duplicating a slide -- SPEC §5.4, §5.5, §5.6.

Assertions run against a saved and reopened package (SPEC §10.2). The engine's
own behaviour is covered in `test_core_clone.py`; what is asserted here is the
verb's contract -- placement, identity, and what a user gets in the file.
"""

from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

from pptx_plus._testing import (
    assert_in_package,
    assert_package_integrity,
    assert_parts_disjoint,
    assert_parts_shared,
    roundtrip,
    saved,
)
from pptx_plus.core.oxml import xpath
from pptx_plus.core.sections import section_lst
from pptx_plus.slides import delete_slide, duplicate_slide, move_slide
from pptx_plus.slides.resolve import SlideIndexError, SlideNotFoundError


def _titles(prs: Presentation) -> list[str]:
    return [slide.shapes.title.text for slide in roundtrip(prs).slides]


def _slide_ids(prs: Presentation) -> list[int]:
    return [int(v) for v in xpath(saved(prs).presentation, "./p:sldIdLst/p:sldId/@id")]


@pytest.fixture
def prs(deck) -> Presentation:
    return Presentation(str(deck("simple4")))


# ---------------------------------------------------------------------------
# Placement -- SPEC §5.4
# ---------------------------------------------------------------------------


def test_the_copy_lands_immediately_after_the_source(prs: Presentation) -> None:
    """The default, matching PowerPoint's own Duplicate Slide."""
    duplicate_slide(prs, 1)
    assert _titles(prs) == ["Slide 1", "Slide 2", "Slide 2", "Slide 3", "Slide 4"]


def test_the_deck_grows_by_one(prs: Presentation) -> None:
    duplicate_slide(prs, 0)
    assert len(roundtrip(prs).slides) == 5


def test_the_returned_slide_is_the_copy(prs: Presentation) -> None:
    copy = duplicate_slide(prs, 0)
    assert copy is not prs.slides[0]


def test_the_returned_slide_is_at_the_expected_index(prs: Presentation) -> None:
    copy = duplicate_slide(prs, 1)
    assert prs.slides.index(copy) == 2


def test_an_explicit_to_index_is_honoured(prs: Presentation) -> None:
    duplicate_slide(prs, 0, to_index=3)
    assert _titles(prs) == ["Slide 1", "Slide 2", "Slide 3", "Slide 1", "Slide 4"]


def test_to_index_zero_puts_the_copy_first(prs: Presentation) -> None:
    duplicate_slide(prs, 3, to_index=0)
    assert _titles(prs) == ["Slide 4", "Slide 1", "Slide 2", "Slide 3", "Slide 4"]


def test_a_negative_to_index_appends(prs: Presentation) -> None:
    """`to_index=-1` is how you ask for the end of the deck."""
    duplicate_slide(prs, 0, to_index=-1)
    assert _titles(prs) == ["Slide 1", "Slide 2", "Slide 3", "Slide 4", "Slide 1"]


def test_duplicating_the_last_slide(prs: Presentation) -> None:
    duplicate_slide(prs, 3)
    assert _titles(prs) == ["Slide 1", "Slide 2", "Slide 3", "Slide 4", "Slide 4"]


def test_a_negative_source_index_is_accepted(prs: Presentation) -> None:
    duplicate_slide(prs, -1)
    assert _titles(prs) == ["Slide 1", "Slide 2", "Slide 3", "Slide 4", "Slide 4"]


def test_a_slide_object_may_be_given(prs: Presentation) -> None:
    duplicate_slide(prs, prs.slides[0])
    assert _titles(prs) == ["Slide 1", "Slide 1", "Slide 2", "Slide 3", "Slide 4"]


@pytest.mark.parametrize("to_index", [5, 99, -6])
def test_an_out_of_range_to_index_raises(prs: Presentation, to_index: int) -> None:
    """The range is the deck *with* the copy in it, so 4 is valid and 5 is not."""
    with pytest.raises(SlideIndexError, match="to_index"):
        duplicate_slide(prs, 0, to_index=to_index)


def test_the_end_of_the_grown_deck_is_in_range(prs: Presentation) -> None:
    duplicate_slide(prs, 0, to_index=4)
    assert _titles(prs)[-1] == "Slide 1"


def test_a_rejected_duplicate_changes_nothing(prs: Presentation) -> None:
    """Validated before the clone, so a rejected call leaves no orphan part."""
    before = saved(prs).partnames
    with pytest.raises(SlideIndexError):
        duplicate_slide(prs, 0, to_index=99)
    assert saved(prs).partnames == before


def test_an_out_of_range_source_raises(prs: Presentation) -> None:
    with pytest.raises(SlideIndexError, match="slide index"):
        duplicate_slide(prs, 9)


def test_a_foreign_slide_raises(prs: Presentation, deck) -> None:
    other = Presentation(str(deck("simple")))
    with pytest.raises(SlideNotFoundError):
        duplicate_slide(prs, other.slides[0])


# ---------------------------------------------------------------------------
# Identity -- SPEC §5.6
# ---------------------------------------------------------------------------


def test_the_copy_gets_a_fresh_slide_id(prs: Presentation) -> None:
    before = set(_slide_ids(prs))
    duplicate_slide(prs, 0)
    assert len(set(_slide_ids(prs)) - before) == 1


def test_the_source_keeps_its_slide_id(prs: Presentation) -> None:
    before = set(_slide_ids(prs))
    duplicate_slide(prs, 0)
    assert before <= set(_slide_ids(prs))


def test_slide_ids_stay_unique(prs: Presentation) -> None:
    duplicate_slide(prs, 0)
    ids = _slide_ids(prs)
    assert len(ids) == len(set(ids))


def test_a_slide_id_freed_by_a_delete_is_not_reused(prs: Presentation) -> None:
    """Anything holding a stale id gets a clean miss, not the wrong slide."""
    freed = _slide_ids(prs)[1]
    delete_slide(prs, 1)
    duplicate_slide(prs, 0)
    assert freed not in _slide_ids(prs)


def test_duplicating_twice_yields_two_distinct_slides(prs: Presentation) -> None:
    """Repeatable, not idempotent. SPEC §5.6."""
    first = duplicate_slide(prs, 0)
    second = duplicate_slide(prs, 0)
    assert first is not second


def test_duplicating_twice_leaves_a_valid_package(prs: Presentation) -> None:
    """What must be idempotent is *allocation*: the second cannot collide."""
    duplicate_slide(prs, 0)
    duplicate_slide(prs, 0)
    assert_package_integrity(saved(prs))


def test_duplicating_a_duplicate_works(prs: Presentation) -> None:
    copy = duplicate_slide(prs, 0)
    duplicate_slide(prs, copy)
    assert len(roundtrip(prs).slides) == 6


# ---------------------------------------------------------------------------
# What is shared and what is copied -- I7
# ---------------------------------------------------------------------------


def test_an_image_is_shared_with_the_copy(deck) -> None:
    presentation = Presentation(str(deck("picture")))
    duplicate_slide(presentation, 0)
    pkg = saved(presentation)
    first, second = pkg.slide_partnames
    assert_parts_shared(pkg, first, second, reltypes=[RT.IMAGE])


def test_the_image_part_is_not_duplicated_in_the_package(deck) -> None:
    """Sharing, observable as the file not growing an image."""
    presentation = Presentation(str(deck("picture")))
    duplicate_slide(presentation, 0)
    media = [n for n in saved(presentation).partnames if n.startswith("/ppt/media/")]
    assert len(media) == 1


def test_a_chart_is_copied_not_shared(deck) -> None:
    """Editing the copy's chart must not reach back into the original."""
    presentation = Presentation(str(deck("chart")))
    duplicate_slide(presentation, 0)
    pkg = saved(presentation)
    first, second = pkg.slide_partnames
    assert_parts_disjoint(pkg, first, second, reltypes=[RT.CHART])


def test_the_copied_chart_owns_its_own_workbook(deck) -> None:
    presentation = Presentation(str(deck("chart")))
    duplicate_slide(presentation, 0)
    pkg = saved(presentation)
    charts = sorted(
        rel.partname
        for slide in pkg.slide_partnames
        for rel in pkg.targets(slide, [RT.CHART])
        if rel.partname
    )
    assert_parts_disjoint(pkg, charts[0], charts[1], reltypes=[RT.PACKAGE])


def test_a_notes_slide_is_copied(deck) -> None:
    presentation = Presentation(str(deck("notes")))
    duplicate_slide(presentation, 1)
    pkg = saved(presentation)
    assert_in_package(pkg, "/ppt/notesSlides/notesSlide2.xml")


def test_the_copied_notes_keep_their_text(deck) -> None:
    presentation = Presentation(str(deck("notes")))
    copy = duplicate_slide(presentation, 1)
    assert copy.notes_slide.notes_text_frame.text == "Speaker notes for slide 2."


def test_with_notes_false_drops_the_notes(deck) -> None:
    presentation = Presentation(str(deck("notes")))
    duplicate_slide(presentation, 1, with_notes=False)
    reopened = roundtrip(presentation)
    assert reopened.slides[2].has_notes_slide is False


def test_with_notes_false_leaves_the_source_notes_alone(deck) -> None:
    presentation = Presentation(str(deck("notes")))
    duplicate_slide(presentation, 1, with_notes=False)
    reopened = roundtrip(presentation)
    assert reopened.slides[1].has_notes_slide is True


def test_two_charts_do_not_collide_on_a_part_name(deck) -> None:
    """The reservation-set regression, observed on the saved zip. SPEC §4.6."""
    presentation = Presentation(str(deck("two_charts")))
    duplicate_slide(presentation, 0)
    charts = [n for n in saved(presentation).partnames if n.startswith("/ppt/charts/chart")]
    assert len(charts) == len(set(charts)) == 4


# ---------------------------------------------------------------------------
# The relationship graph stays closed -- I4
# ---------------------------------------------------------------------------


@pytest.mark.parametrize(
    "name", ["simple", "picture", "shared_picture", "notes", "chart", "two_charts", "hyperlink"]
)
def test_duplicating_leaves_a_valid_package(deck, name: str) -> None:
    presentation = Presentation(str(deck(name)))
    duplicate_slide(presentation, 0)
    assert_package_integrity(saved(presentation))


def test_the_gap_rids_deck_duplicates_correctly(deck) -> None:
    """The deck built to make a skipped relationship rewrite fail loudly."""
    presentation = Presentation(str(deck("gap_rids")))
    duplicate_slide(presentation, 0)
    assert_package_integrity(saved(presentation))


def test_a_duplicated_deck_survives_a_second_round_trip(prs: Presentation) -> None:
    duplicate_slide(prs, 0)
    assert_package_integrity(saved(roundtrip(prs)))


def test_the_copy_renders_the_same_shapes(deck) -> None:
    """Fidelity, at the level a test can actually check."""
    presentation = Presentation(str(deck("picture")))
    copy = duplicate_slide(presentation, 0)
    assert len(copy.shapes) == len(presentation.slides[0].shapes)


def test_a_hyperlink_survives_the_copy(deck) -> None:
    presentation = Presentation(str(deck("hyperlink")))
    duplicate_slide(presentation, 0)
    pkg = saved(presentation)
    external = [
        rel.target
        for slide in pkg.slide_partnames
        for rel in pkg.rels(slide).values()
        if rel.is_external
    ]
    assert len(external) == 2


def test_a_slide_jump_still_points_at_the_original(deck) -> None:
    """The copy links to slide 1; it does not bring a second slide 1 along."""
    presentation = Presentation(str(deck("slide_jump")))
    duplicate_slide(presentation, 1)
    pkg = saved(presentation)
    assert len(pkg.slide_partnames) == 3
    targets = {
        rel.partname for slide in pkg.slide_partnames for rel in pkg.targets(slide, [RT.SLIDE])
    }
    assert targets == {"/ppt/slides/slide1.xml"}


# ---------------------------------------------------------------------------
# Sections
# ---------------------------------------------------------------------------


def _sections(prs: Presentation) -> list[list[int]]:
    root = section_lst(prs)
    if root is None:
        return []
    return [
        [int(value) for value in xpath(id_list, "./p14:sldId/@id")]
        for id_list in xpath(root, "./p14:section/p14:sldIdLst")
    ]


def test_the_copy_joins_its_sources_section(deck) -> None:
    """Matching PowerPoint: a duplicate lands in the section it came from."""
    presentation = Presentation(str(deck("sections")))
    duplicate_slide(presentation, 0)
    assert _sections(presentation) == [[256, 260, 257], [258, 259]]


def test_a_copy_at_a_section_boundary_stays_in_the_earlier_section(deck) -> None:
    """Duplicating the last slide of section 1 must not open section 2."""
    presentation = Presentation(str(deck("sections")))
    duplicate_slide(presentation, 1)
    assert _sections(presentation) == [[256, 257, 260], [258, 259]]


def test_a_copy_placed_elsewhere_joins_that_section(deck) -> None:
    presentation = Presentation(str(deck("sections")))
    duplicate_slide(presentation, 0, to_index=-1)
    assert _sections(presentation) == [[256, 257], [258, 259, 260]]


def test_sections_still_cover_the_grown_deck(deck) -> None:
    presentation = Presentation(str(deck("sections")))
    duplicate_slide(presentation, 0)
    flattened = [sld_id for section in _sections(presentation) for sld_id in section]
    assert flattened == _slide_ids(presentation)


def test_a_sectioned_duplicate_is_valid(deck) -> None:
    presentation = Presentation(str(deck("sections")))
    duplicate_slide(presentation, 2)
    assert_package_integrity(saved(presentation))


def test_a_deck_without_sections_gains_none(prs: Presentation) -> None:
    duplicate_slide(prs, 0)
    assert section_lst(roundtrip(prs)) is None


def test_the_copy_is_not_added_to_any_custom_show(deck) -> None:
    """A custom show is a curated list; a duplicate does not join it uninvited."""
    presentation = Presentation(str(deck("custom_show")))
    duplicate_slide(presentation, 0)
    shows = xpath(saved(presentation).presentation, "./p:custShowLst//p:sld")
    assert len(shows) == 2


# ---------------------------------------------------------------------------
# Composition with the other verbs
# ---------------------------------------------------------------------------


def test_duplicate_then_delete_the_original(prs: Presentation) -> None:
    duplicate_slide(prs, 0)
    delete_slide(prs, 0)
    assert _titles(prs) == ["Slide 1", "Slide 2", "Slide 3", "Slide 4"]


def test_duplicate_then_delete_leaves_the_shared_image(deck) -> None:
    presentation = Presentation(str(deck("picture")))
    duplicate_slide(presentation, 0)
    delete_slide(presentation, 0)
    assert_package_integrity(saved(presentation))
    assert_in_package(saved(presentation), "/ppt/media/image1.gif")


def test_duplicate_then_move(prs: Presentation) -> None:
    duplicate_slide(prs, 0)
    move_slide(prs, 1, 4)
    assert _titles(prs) == ["Slide 1", "Slide 2", "Slide 3", "Slide 4", "Slide 1"]


def test_deleting_a_copy_does_not_disturb_the_original(deck) -> None:
    """The copy owns its chart, so collecting it must not take the source's."""
    presentation = Presentation(str(deck("chart")))
    copy = duplicate_slide(presentation, 0)
    delete_slide(presentation, copy)
    pkg = saved(presentation)
    assert_package_integrity(pkg)
    assert len(pkg.targets(pkg.slide_partnames[0], [RT.CHART])) == 1

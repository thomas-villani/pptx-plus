"""Part-graph classification and cloning -- SPEC §4.5.

The engine's behaviour is stated here in terms of parts and relationships;
`test_slides_duplicate.py` asserts the same properties through the verb.
"""

from __future__ import annotations

from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

from pptx_plus.core.clone import ClonePolicy, clone_part_graph
from pptx_plus.core.ns import qn
from pptx_plus.core.oxml import part_root, xpath
from pptx_plus.core.partgraph import Disposition, RelEdge, classify, rel_edges
from pptx_plus.core.reltypes import RT_DIAGRAM_DRAWING


def _clone_first_slide(prs: Presentation, **kwargs: object) -> object:
    part = prs.slides[0].part
    policy = ClonePolicy(**kwargs) if kwargs else None  # type: ignore[arg-type]
    return clone_part_graph(part, into=part.package, policy=policy)


def _edge(prs: Presentation, reltype: str, index: int = 0):
    return next(e for e in rel_edges(prs.slides[index].part) if e.reltype == reltype)


# ---------------------------------------------------------------------------
# classify
# ---------------------------------------------------------------------------


def test_a_layout_is_structural(deck) -> None:
    """Two slides sharing a layout is correct; duplicating it would not be."""
    prs = Presentation(str(deck("simple")))
    assert classify(_edge(prs, RT.SLIDE_LAYOUT)) is Disposition.STRUCTURAL


def test_an_image_is_shared(deck) -> None:
    prs = Presentation(str(deck("picture")))
    assert classify(_edge(prs, RT.IMAGE)) is Disposition.SHARE


def test_a_chart_is_deep_cloned(deck) -> None:
    """No special case: RT.CHART is in no reuse set, so the general rule copies it."""
    prs = Presentation(str(deck("chart")))
    assert classify(_edge(prs, RT.CHART)) is Disposition.DEEP


def test_an_external_relationship_is_external(deck) -> None:
    prs = Presentation(str(deck("hyperlink")))
    assert classify(_edge(prs, RT.HYPERLINK)) is Disposition.EXTERNAL


def test_a_notes_slide_is_deep_cloned(deck) -> None:
    prs = Presentation(str(deck("notes")))
    assert classify(_edge(prs, RT.NOTES_SLIDE, index=1)) is Disposition.DEEP


def test_a_slide_reference_is_reused(deck) -> None:
    """A slide-jump target is another slide, not this slide's content."""
    prs = Presentation(str(deck("slide_jump")))
    assert classify(_edge(prs, RT.SLIDE, index=1)) is Disposition.REUSE


def test_a_shareable_type_with_its_own_relationships_is_deep_cloned(deck) -> None:
    """The safety net clause, and the half that does the real work.

    Membership in SHARE_RELTYPES is necessary but not sufficient: anything
    owning a sub-graph is a mutable unit whatever the reltype table says. If a
    future Office release ships an image part with a sub-relationship, it gets
    copied rather than silently aliased.
    """
    prs = Presentation(str(deck("chart")))
    chart_part = _edge(prs, RT.CHART).target_part
    assert chart_part is not None and len(chart_part.rels) > 0
    disguised = RelEdge("rId9", RT.IMAGE, False, chart_part, "")
    assert classify(disguised) is Disposition.DEEP


def test_an_embedded_workbook_is_deep_cloned(deck) -> None:
    """RT.PACKAGE is not shareable, so the chart's workbook is copied."""
    prs = Presentation(str(deck("chart")))
    chart_part = _edge(prs, RT.CHART).target_part
    assert chart_part is not None
    edge = next(e for e in rel_edges(chart_part) if e.reltype == RT.PACKAGE)
    assert classify(edge) is Disposition.DEEP


# ---------------------------------------------------------------------------
# rel_edges
# ---------------------------------------------------------------------------


def test_rel_edges_snapshots_every_relationship(deck) -> None:
    prs = Presentation(str(deck("picture")))
    part = prs.slides[0].part
    assert len(rel_edges(part)) == len(part.rels)


def test_rel_edges_carries_the_target_ref_for_an_external(deck) -> None:
    prs = Presentation(str(deck("hyperlink")))
    assert _edge(prs, RT.HYPERLINK).target_ref.startswith("https://")


def test_rel_edges_has_no_target_part_for_an_external(deck) -> None:
    prs = Presentation(str(deck("hyperlink")))
    assert _edge(prs, RT.HYPERLINK).target_part is None


def test_reading_relationships_creates_no_notes_slide(deck) -> None:
    """SPEC §9.8: `notes_slide` is a lazy property that *creates* parts.

    Touching it on an inspection path would create a notes slide, a notes
    master, and a theme -- silently growing a package this was only meant to
    read.
    """
    prs = Presentation(str(deck("simple")))
    part = prs.slides[0].part
    rel_edges(part)
    assert part.has_notes_slide is False


# ---------------------------------------------------------------------------
# clone_part_graph
# ---------------------------------------------------------------------------


def test_the_clone_is_a_new_part(deck) -> None:
    prs = Presentation(str(deck("simple")))
    assert _clone_first_slide(prs).root is not prs.slides[0].part


def test_the_clone_gets_a_fresh_partname(deck) -> None:
    prs = Presentation(str(deck("simple")))
    assert _clone_first_slide(prs).root.partname != prs.slides[0].part.partname


def test_the_clone_reuses_the_layout(deck) -> None:
    prs = Presentation(str(deck("simple")))
    source = prs.slides[0].part
    clone = _clone_first_slide(prs).root
    assert _target(clone, RT.SLIDE_LAYOUT) is _target(source, RT.SLIDE_LAYOUT)


def test_the_clone_shares_the_image(deck) -> None:
    prs = Presentation(str(deck("picture")))
    source = prs.slides[0].part
    clone = _clone_first_slide(prs).root
    assert _target(clone, RT.IMAGE) is _target(source, RT.IMAGE)


def test_the_clone_copies_the_chart(deck) -> None:
    prs = Presentation(str(deck("chart")))
    source = prs.slides[0].part
    clone = _clone_first_slide(prs).root
    assert _target(clone, RT.CHART) is not _target(source, RT.CHART)


def test_the_cloned_chart_gets_its_own_workbook(deck) -> None:
    """Two levels down: cloning is recursive, not one level deep."""
    prs = Presentation(str(deck("chart")))
    source_chart = _target(prs.slides[0].part, RT.CHART)
    clone_chart = _target(_clone_first_slide(prs).root, RT.CHART)
    assert _target(clone_chart, RT.PACKAGE) is not _target(source_chart, RT.PACKAGE)


def test_the_cloned_workbook_is_byte_identical(deck) -> None:
    """Copied as bytes, never reparsed. SPEC §8.1."""
    prs = Presentation(str(deck("chart")))
    source_chart = _target(prs.slides[0].part, RT.CHART)
    clone_chart = _target(_clone_first_slide(prs).root, RT.CHART)
    assert _target(clone_chart, RT.PACKAGE).blob == _target(source_chart, RT.PACKAGE).blob


def test_two_charts_get_distinct_partnames(deck) -> None:
    """The reservation-set regression. SPEC §4.6.

    `next_partname` derives its used-name set from a walk of the relationship
    graph, so a part constructed but not yet attached is invisible to it.
    Without the reservation set both cloned charts are handed the same name
    and one silently overwrites the other in the zip.
    """
    prs = Presentation(str(deck("two_charts")))
    clone = _clone_first_slide(prs).root
    names = {
        str(rel.target_part.partname) for rel in clone.rels.values() if rel.reltype == RT.CHART
    }
    assert len(names) == 2


def test_an_external_relationship_is_re_minted(deck) -> None:
    prs = Presentation(str(deck("hyperlink")))
    clone = _clone_first_slide(prs).root
    external = [rel for rel in clone.rels.values() if rel.is_external]
    assert [rel.target_ref for rel in external] == ["https://github.com/scanny/python-pptx"]


# ---------------------------------------------------------------------------
# The clone map -- cycles, notes, and slide references
# ---------------------------------------------------------------------------


def test_a_notes_slide_is_cloned(deck) -> None:
    prs = Presentation(str(deck("notes")))
    part = prs.slides[1].part
    clone = clone_part_graph(part, into=part.package).root
    assert _target(clone, RT.NOTES_SLIDE) is not _target(part, RT.NOTES_SLIDE)


def test_the_cloned_notes_slide_points_back_at_the_clone(deck) -> None:
    """The clone-map invariant, and the cycle it resolves. SPEC §4.5.

    A NotesSlidePart holds a relationship *back* to its slide. By the time it
    is cloned the source slide is already in the map, so the back-reference
    lands on the clone -- for free, with no special case for notes anywhere in
    the engine.
    """
    prs = Presentation(str(deck("notes")))
    part = prs.slides[1].part
    clone = clone_part_graph(part, into=part.package).root
    assert _target(_target(clone, RT.NOTES_SLIDE), RT.SLIDE) is clone


def test_with_notes_false_omits_the_notes_slide(deck) -> None:
    prs = Presentation(str(deck("notes")))
    part = prs.slides[1].part
    clone = clone_part_graph(part, into=part.package, policy=ClonePolicy(with_notes=False)).root
    assert [rel for rel in clone.rels.values() if rel.reltype == RT.NOTES_SLIDE] == []


def test_cloning_a_notes_slide_terminates(deck) -> None:
    """slide -> notesSlide -> slide is a real cycle; mapping before recursing
    is what stops it recursing forever."""
    prs = Presentation(str(deck("notes")))
    part = prs.slides[1].part
    assert clone_part_graph(part, into=part.package).root is not None


def test_a_slide_jump_target_is_not_duplicated(deck) -> None:
    """The copy links to slide 1; it does not bring a second slide 1 with it."""
    prs = Presentation(str(deck("slide_jump")))
    part = prs.slides[1].part
    clone = clone_part_graph(part, into=part.package).root
    assert _target(clone, RT.SLIDE) is _target(part, RT.SLIDE)


# ---------------------------------------------------------------------------
# The relationship-id rewrite, at the part level
# ---------------------------------------------------------------------------


def test_the_clone_is_not_attached_to_anything(deck) -> None:
    """Minting the relationship that puts it in the deck is the caller's job."""
    prs = Presentation(str(deck("simple")))
    clone = _clone_first_slide(prs).root
    assert not any(
        rel.target_part is clone for rel in prs.part.rels.values() if not rel.is_external
    )


def test_every_rel_id_in_the_clone_resolves(deck) -> None:
    prs = Presentation(str(deck("picture")))
    clone = _clone_first_slide(prs).root
    values = {
        str(v) for v in xpath(part_root(clone), "//@*[namespace-uri()=$ns]", ns=qn("r:id")[1:-3])
    }
    assert values <= set(clone.rels)


def test_the_gap_rids_clone_rewrites_the_embed(deck) -> None:
    """**The primary regression for the whole library.**

    The source's relationships are {rId1, rId3} -- the ordinary result of
    deleting a shape in PowerPoint. `_next_rId` fills gaps, so the clone's are
    {rId1, rId2}, and an unrewritten `r:embed="rId3"` would dangle.

    In a freshly authored deck the clone's ids come out identical to the
    source's and a skipped rewrite passes by accident. This deck removes the
    accident, which is the only reason it exists.
    """
    prs = Presentation(str(deck("gap_rids")))
    source = prs.slides[0].part
    clone = _clone_first_slide(prs).root

    assert xpath(part_root(source), "//@r:embed") == ["rId3"]
    assert xpath(part_root(clone), "//@r:embed") == ["rId2"]


def test_the_rewritten_embed_resolves_to_the_same_image(deck) -> None:
    """Rewritten to a *correct* id, not merely to a different one."""
    prs = Presentation(str(deck("gap_rids")))
    source = prs.slides[0].part
    clone = _clone_first_slide(prs).root
    assert clone.rels["rId2"].target_part is source.rels["rId3"].target_part


def test_a_chart_sub_part_is_rewritten_too(deck) -> None:
    """The pass runs per cloned XML part, not once over the slide. SPEC §4.4."""
    prs = Presentation(str(deck("chart")))
    clone_chart = _target(_clone_first_slide(prs).root, RT.CHART)
    external = xpath(part_root(clone_chart), "./c:externalData/@r:id")
    assert [str(v) for v in external] == [
        rel.rId for rel in clone_chart.rels.values() if rel.reltype == RT.PACKAGE
    ]


def test_diagram_reltypes_are_all_deep(deck) -> None:
    """Asserted from the constant, so a test list cannot drift from the code."""
    for reltype in (
        RT.DIAGRAM_DATA,
        RT.DIAGRAM_LAYOUT,
        RT.DIAGRAM_COLORS,
        RT.DIAGRAM_QUICK_STYLE,
        RT_DIAGRAM_DRAWING,
    ):
        assert classify(RelEdge("rId9", reltype, False, _blob_part(deck), "")) is Disposition.DEEP


# ---------------------------------------------------------------------------
# CloneResult
# ---------------------------------------------------------------------------


def test_the_result_reports_the_parts_it_created(deck) -> None:
    prs = Presentation(str(deck("chart")))
    result = _clone_first_slide(prs)
    assert len(result.cloned_parts) == 3  # slide, chart, workbook


def test_the_result_reports_what_it_shared(deck) -> None:
    prs = Presentation(str(deck("picture")))
    result = _clone_first_slide(prs)
    assert [p.content_type for p in result.shared] == ["image/gif"]


def test_the_result_maps_source_parts_to_clones(deck) -> None:
    prs = Presentation(str(deck("simple")))
    source = prs.slides[0].part
    result = clone_part_graph(source, into=source.package)
    assert result.part_map[id(source)] is result.root


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _target(part: object, reltype: str) -> object:
    return next(rel.target_part for rel in part.rels.values() if rel.reltype == reltype)


def _blob_part(deck) -> object:
    """A part with no relationships of its own, for synthetic classify() edges."""
    prs = Presentation(str(deck("picture")))
    return next(
        rel.target_part for rel in prs.slides[0].part.rels.values() if rel.reltype == RT.IMAGE
    )

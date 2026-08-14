"""Element construction and query chokepoint — SPEC §4.2, §9.2."""

from __future__ import annotations

import pytest
from lxml import etree
from pptx import Presentation

from pptx_plus.core.ns import InvalidNamespaceError, qn
from pptx_plus.core.oxml import (
    el,
    ordered_insert,
    part_root,
    remove,
    sld_id_lst,
    sub,
    xpath,
)

# --------------------------------------------------------------------------
# el / sub
# --------------------------------------------------------------------------


def test_el_builds_a_namespaced_element() -> None:
    assert el("p:sldId").tag == qn("p:sldId")


def test_el_sets_a_namespaced_attribute() -> None:
    node = el("p:sldId", **{"r:id": "rId3"})
    assert node.get(qn("r:id")) == "rId3"


def test_el_sets_a_bare_attribute() -> None:
    """p:sldId carries a namespaced r:id next to a bare id; both are ordinary."""
    assert el("p:sldId", id="256").get("id") == "256"


def test_el_rejects_an_unknown_prefix() -> None:
    with pytest.raises(InvalidNamespaceError):
        el("nope:thing")


def test_el_rejects_an_unknown_attribute_prefix() -> None:
    with pytest.raises(InvalidNamespaceError):
        el("p:sldId", **{"nope:x": "1"})


def test_el_declares_only_its_own_prefix_outside_the_build_map() -> None:
    """An extension part should not carry irrelevant namespace declarations."""
    assert set(el("p14:media").nsmap) == {"p14"}


def test_sub_appends_to_the_parent() -> None:
    parent = el("p:sldIdLst")
    child = sub(parent, "p:sldId", id="256")
    assert list(parent) == [child]


# --------------------------------------------------------------------------
# xpath
# --------------------------------------------------------------------------


def test_xpath_resolves_prefixes_from_nsmap() -> None:
    parent = el("p:sldIdLst")
    sub(parent, "p:sldId", id="256")
    assert len(xpath(parent, "./p:sldId")) == 1


def test_xpath_binds_variables() -> None:
    parent = el("p:sldIdLst")
    sub(parent, "p:sldId", **{"id": "256", "r:id": "rId2"})
    sub(parent, "p:sldId", **{"id": "257", "r:id": "rId3"})
    found = xpath(parent, "./p:sldId[@r:id=$rid]", rid="rId3")
    assert found[0].get("id") == "257"


def test_xpath_variable_binding_survives_a_quote() -> None:
    """The concrete reason variables beat f-string interpolation.

    A value containing a quote would terminate the literal early and produce
    either a syntax error or, worse, a silently different expression.
    """
    parent = el("p:sldIdLst")
    sub(parent, "p:sldId", **{"id": "it's", "r:id": "rId2"})
    assert len(xpath(parent, "./p:sldId[@id=$v]", v="it's")) == 1


def test_xpath_wraps_a_scalar_result_in_a_list() -> None:
    assert xpath(el("p:sldIdLst"), "count(./p:sldId)") == [0.0]


def test_xpath_reaches_an_extension_namespace() -> None:
    """python-pptx's own xpath() cannot do this — see test_core_ns.py."""
    root = el("p:extLst")
    sub(root, "mc:AlternateContent")
    assert len(xpath(root, "./mc:AlternateContent")) == 1


# --------------------------------------------------------------------------
# remove
# --------------------------------------------------------------------------


def test_remove_detaches_from_the_parent() -> None:
    parent = el("p:sldIdLst")
    child = sub(parent, "p:sldId")
    remove(child)
    assert len(parent) == 0


def test_remove_is_a_noop_on_a_detached_node() -> None:
    """Makes the delete paths re-entrant rather than order-sensitive."""
    remove(el("p:sldId"))  # must not raise


# --------------------------------------------------------------------------
# ordered_insert
# --------------------------------------------------------------------------

_ORDER = ("p:sldMasterIdLst", "p:sldIdLst", "p:sldSz", "p:notesSz")


def _tags(parent: etree._Element) -> list[str]:
    return [etree.QName(child).localname for child in parent]


def test_ordered_insert_places_before_a_later_sibling() -> None:
    parent = el("p:presentation")
    sub(parent, "p:sldSz")
    ordered_insert(parent, el("p:sldIdLst"), _ORDER)
    assert _tags(parent) == ["sldIdLst", "sldSz"]


def test_ordered_insert_appends_when_no_later_sibling_exists() -> None:
    parent = el("p:presentation")
    sub(parent, "p:sldMasterIdLst")
    ordered_insert(parent, el("p:sldIdLst"), _ORDER)
    assert _tags(parent) == ["sldMasterIdLst", "sldIdLst"]


def test_ordered_insert_is_idempotent() -> None:
    """Calling twice leaves one element, not two."""
    parent = el("p:presentation")
    sub(parent, "p:sldSz")
    ordered_insert(parent, el("p:sldIdLst"), _ORDER)
    ordered_insert(parent, el("p:sldIdLst"), _ORDER)
    assert _tags(parent) == ["sldIdLst", "sldSz"]


def test_ordered_insert_rejects_a_tag_absent_from_the_order() -> None:
    with pytest.raises(ValueError, match="not in the declared child order"):
        ordered_insert(el("p:presentation"), el("p:sldId"), _ORDER)


# --------------------------------------------------------------------------
# The private-API quarantine — SPEC §14.2
# --------------------------------------------------------------------------


def test_part_root_returns_the_live_tree() -> None:
    """Live, not a copy: mutating the returned element mutates the part."""
    prs = Presentation()
    root = part_root(prs.part)
    assert root is part_root(prs.part)


def test_sld_id_lst_returns_the_slide_list() -> None:
    prs = Presentation()
    assert etree.QName(sld_id_lst(prs)).localname == "sldIdLst"


def test_sld_id_lst_agrees_with_python_pptx() -> None:
    """The accessor must find the same element the naive recipe reaches for.

    `prs.slides._sldIdLst` is where every StackOverflow recipe starts. Going
    through `core` instead only helps if it lands on the same element.
    """
    prs = Presentation()
    assert sld_id_lst(prs) is prs.slides._sldIdLst


def test_sld_id_lst_tracks_added_slides() -> None:
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    assert len(xpath(sld_id_lst(prs), "./p:sldId")) == 1

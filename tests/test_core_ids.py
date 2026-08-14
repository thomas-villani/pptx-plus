"""Slide-id allocation and range validation — SPEC §4, §5.6."""

from __future__ import annotations

import pytest
from pptx import Presentation

from pptx_plus.core.ids import (
    MAX_SLIDE_ID,
    MIN_SLIDE_ID,
    SlideIdRangeError,
    next_slide_id,
    used_slide_ids,
    validate_slide_id,
)
from pptx_plus.core.oxml import el, sld_id_lst, sub


def _lst(*ids: int) -> object:
    lst = el("p:sldIdLst")
    for index, value in enumerate(ids):
        sub(lst, "p:sldId", **{"id": str(value), "r:id": f"rId{index + 2}"})
    return lst


# --------------------------------------------------------------------------
# validate_slide_id
# --------------------------------------------------------------------------


@pytest.mark.parametrize("value", [MIN_SLIDE_ID, 1000, MAX_SLIDE_ID])
def test_validate_accepts_the_legal_range(value: int) -> None:
    assert validate_slide_id(value) == value


@pytest.mark.parametrize("value", [0, 255, MAX_SLIDE_ID + 1, -1])
def test_validate_rejects_outside_the_legal_range(value: int) -> None:
    with pytest.raises(SlideIdRangeError, match="ST_SlideId"):
        validate_slide_id(value)


def test_slide_id_range_error_is_a_value_error() -> None:
    assert issubclass(SlideIdRangeError, ValueError)


# --------------------------------------------------------------------------
# used_slide_ids
# --------------------------------------------------------------------------


def test_used_ids_are_returned_in_document_order() -> None:
    assert used_slide_ids(_lst(300, 256, 999)) == [300, 256, 999]


def test_used_ids_preserves_duplicates() -> None:
    """A duplicate id is a defect worth reporting, so it is not collapsed here."""
    assert used_slide_ids(_lst(256, 256)) == [256, 256]


def test_used_ids_of_an_empty_list() -> None:
    assert used_slide_ids(_lst()) == []


# --------------------------------------------------------------------------
# next_slide_id
# --------------------------------------------------------------------------


def test_next_id_of_an_empty_deck_is_the_minimum() -> None:
    assert next_slide_id(_lst()) == MIN_SLIDE_ID


def test_next_id_is_one_past_the_maximum_used() -> None:
    assert next_slide_id(_lst(256, 257)) == 258


def test_next_id_does_not_reuse_a_gap() -> None:
    """SPEC §5.6: slide ids are never reused.

    This is the contract that makes a stale slide id fail cleanly instead of
    silently resolving to a different slide, so it is asserted directly rather
    than left as an implementation detail of `add_sldId`.
    """
    assert next_slide_id(_lst(256, 258)) == 259


def test_next_id_raises_when_the_space_is_exhausted() -> None:
    with pytest.raises(SlideIdRangeError, match="cannot allocate"):
        next_slide_id(_lst(MAX_SLIDE_ID))


def test_next_id_agrees_with_python_pptx() -> None:
    """Predicting the id must match what `add_sldId` actually allocates.

    `next_slide_id` exists so tests and callers can ask "what would come next"
    without mutating anything. That is only useful if it agrees with the
    allocator it mirrors.
    """
    prs = Presentation()
    for _ in range(3):
        prs.slides.add_slide(prs.slide_layouts[6])
    lst = sld_id_lst(prs)
    predicted = next_slide_id(lst)

    prs.slides.add_slide(prs.slide_layouts[6])
    assert used_slide_ids(lst)[-1] == predicted


def test_allocated_ids_are_in_range() -> None:
    prs = Presentation()
    for _ in range(3):
        prs.slides.add_slide(prs.slide_layouts[6])
    assert all(validate_slide_id(value) for value in used_slide_ids(sld_id_lst(prs)))

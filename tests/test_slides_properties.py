"""Property tests over sequences of slide operations.

The one place `hypothesis` earns its dependency. Index arithmetic across three
verbs -- with negative indices, two different length bases, and `to_index`
meaning a position in the *resulting* deck -- is exactly the bug class
hand-written examples miss, because the examples are written by the same person
who wrote the arithmetic. When it does fail, the shrinker hands back a
two-call reproduction rather than a forty-step trace.

Every sequence ends with the full integrity battery on the saved package.
"""

from __future__ import annotations

from dataclasses import dataclass

import pytest
from hypothesis import HealthCheck, given, settings
from hypothesis import strategies as st
from pptx import Presentation

from pptx_plus._testing import assert_package_integrity, roundtrip, saved
from pptx_plus.core.oxml import xpath
from pptx_plus.core.sections import section_lst
from pptx_plus.slides import delete_slide, duplicate_slide, move_slide


@dataclass(frozen=True)
class Op:
    """One operation, resolved against the deck length at the time it runs."""

    verb: str
    #: Fractions of the current deck length, so an operation is always in
    #: range however the preceding ones changed it. Generating raw indices
    #: would make almost every sequence a chain of SlideIndexErrors, and the
    #: interesting failures are in the arithmetic, not the guards.
    source: float
    target: float


OPS = st.builds(
    Op,
    verb=st.sampled_from(["delete", "move", "duplicate"]),
    source=st.floats(min_value=0, max_value=0.999),
    target=st.floats(min_value=0, max_value=0.999),
)


def _apply(prs: Presentation, op: Op) -> None:
    count = len(prs.slides)
    if count == 0:
        return
    source = int(op.source * count)

    if op.verb == "delete":
        delete_slide(prs, source)
    elif op.verb == "move":
        move_slide(prs, source, int(op.target * count))
    else:
        duplicate_slide(prs, source, to_index=int(op.target * (count + 1)))


def _titles(prs: Presentation) -> list[str]:
    return [slide.shapes.title.text for slide in prs.slides]


def _sections(prs: Presentation) -> list[list[int]]:
    root = section_lst(prs)
    if root is None:
        return []
    return [
        [int(value) for value in xpath(id_list, "./p14:sldId/@id")]
        for id_list in xpath(root, "./p14:section/p14:sldIdLst")
    ]


# `deck` is session-scoped and only read from, so the health check about
# function-scoped fixtures does not apply; the decks are never mutated in
# place, only opened.
SETTINGS = settings(
    max_examples=60,
    deadline=None,
    suppress_health_check=[HealthCheck.function_scoped_fixture],
)


@given(ops=st.lists(OPS, min_size=1, max_size=8))
@SETTINGS
def test_any_sequence_leaves_a_valid_package(deck, ops: list[Op]) -> None:
    """The battery holds however the three verbs are composed."""
    prs = Presentation(str(deck("simple4")))
    for op in ops:
        _apply(prs, op)
    assert_package_integrity(saved(prs))


@given(ops=st.lists(OPS, min_size=1, max_size=8))
@SETTINGS
def test_the_running_order_survives_a_round_trip(deck, ops: list[Op]) -> None:
    """What the in-memory deck says and what the file says must agree.

    They are two different objects (SPEC §3.5), and every other test here
    asserts on one or the other. This one asserts they have not diverged.
    """
    prs = Presentation(str(deck("simple4")))
    for op in ops:
        _apply(prs, op)
    assert _titles(roundtrip(prs)) == _titles(prs)


@given(ops=st.lists(OPS, min_size=1, max_size=6))
@SETTINGS
def test_slide_ids_stay_unique(deck, ops: list[Op]) -> None:
    """Duplication allocates, deletion frees, and neither may collide."""
    prs = Presentation(str(deck("simple4")))
    for op in ops:
        _apply(prs, op)
    ids = [int(v) for v in xpath(saved(prs).presentation, "./p:sldIdLst/p:sldId/@id")]
    assert len(ids) == len(set(ids))


@given(ops=st.lists(OPS, min_size=1, max_size=6))
@SETTINGS
def test_sections_always_partition_the_deck(deck, ops: list[Op]) -> None:
    """The flattened sections equal the running order, after any sequence.

    The invariant that makes sections a partition rather than an arbitrary
    grouping, and the one a hand-written test is least likely to break --
    every interesting violation needs a slide crossing a boundary that an
    earlier operation moved.
    """
    prs = Presentation(str(deck("sections")))
    for op in ops:
        _apply(prs, op)
    pkg = saved(prs)
    deck_ids = [int(v) for v in xpath(pkg.presentation, "./p:sldIdLst/p:sldId/@id")]
    assert [sid for section in _sections(prs) for sid in section] == deck_ids


@given(count=st.integers(min_value=0, max_value=4))
@SETTINGS
def test_deleting_every_slide_is_always_legal(deck, count: int) -> None:
    prs = Presentation(str(deck("simple4")))
    for _ in range(count):
        delete_slide(prs, 0)
    assert_package_integrity(saved(prs))


@given(index=st.integers(min_value=0, max_value=3), to_index=st.integers(min_value=0, max_value=3))
@SETTINGS
def test_a_move_is_always_reversible(deck, index: int, to_index: int) -> None:
    """Moving there and back is the identity, for every pair of positions."""
    prs = Presentation(str(deck("simple4")))
    before = _titles(prs)
    move_slide(prs, index, to_index)
    move_slide(prs, to_index, index)
    assert _titles(prs) == before


@given(index=st.integers(min_value=0, max_value=3))
@SETTINGS
def test_duplicate_then_delete_the_copy_is_the_identity(deck, index: int) -> None:
    """The copy is separable: removing it restores the deck exactly."""
    prs = Presentation(str(deck("simple4")))
    before = _titles(prs)
    copy = duplicate_slide(prs, index)
    delete_slide(prs, copy)
    assert _titles(prs) == before


@pytest.mark.parametrize("name", ["picture", "chart", "notes"])
@given(ops=st.lists(OPS, min_size=1, max_size=4))
@SETTINGS
def test_media_bearing_decks_survive_any_sequence(deck, name: str, ops: list[Op]) -> None:
    """Shared images and owned charts, under composition rather than in isolation."""
    prs = Presentation(str(deck(name)))
    for op in ops:
        _apply(prs, op)
    assert_package_integrity(saved(prs))
